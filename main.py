import os
import sqlite3
import time # <--- أضف هذا السطر
from datetime import date
from datetime import datetime, timedelta
import telebot
from telebot.types import ChatPermissions
import threading
import logging
from dotenv import load_dotenv
from telebot.types import InlineKeyboardMarkup, InlineKeyboardButton
import docx
import fitz                     # PyMuPDF
import google.generativeai as genai
import requests
import cohere
from groq import Groq
import json
import re
from pptx import Presentation
import traceback
import threading
import queue
import time
import random
from flask import Flask, render_template


# متغيرات البيئة
load_dotenv()
BOT_TOKEN = os.getenv("BOT_TOKEN")
GROUP_ID = int(os.getenv("GROUP_ID"))
ADMIN_ID = int(os.getenv("ADMIN_ID"))
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
WEBHOOK_URL = os.getenv('WEBHOOK_URL')
allowed_channels = set()
env_channels = os.getenv("ALLOWED_CHANNELS", "")
if env_channels.strip():
    allowed_channels = set(map(int, env_channels.split(",")))
BOT_TOKEN_2 = os.getenv("BOT_TOKEN_2")
BOT_TOKEN_3 = os.getenv("BOT_TOKEN_3")

OCR_API_KEY = os.getenv("OCR_SPACE_API_KEY", "helloworld")  # مفتاح افتراضي للاختبار
bot = telebot.TeleBot(BOT_TOKEN)
bot2 = telebot.TeleBot(BOT_TOKEN_2)
bot3 = telebot.TeleBot(BOT_TOKEN_3)




# -------------------------------------------------------------------
# --------- Notofication code section ---
# -------------------------------------------------------------------

def notify_admin(action: str, username: str, user_id: int):
    """
    يرسل إشعار للإدمن بخصوص عملية معينة.
    action: نوع العملية (مثلاً: "توليد أنكي آلي")
    """
    bot3.send_message(
        ADMIN_ID,
        f"نجاح✔️، عملية {action} ل: {username} | UID: {user_id}"
    )


def notify_process_info(uid, file_id, username):
    try:

        # إرسال إشعار إلى الأدمن فقط
        # بدل user_id استخدم uid
        bot3.send_message(
            ADMIN_ID,
            f"📁 تم الانتهاء من معالجة ملف من المستخدم:\n"
            f"👤 ID: {uid}\n"
            f"🔗 File ID: {file_id or 'no-file'}\n"
            f"💬 Username: @{username if username else 'N/A'}"
            )
    except Exception:
        logging.exception("notify admin failed")


from datetime import datetime, timedelta

usage_count = {}
last_feedback_time = {}

def maybe_send_feedback_request(uid: int, chat_id: int):
    now = datetime.utcnow()
    usage_count[uid] = usage_count.get(uid, 0) + 1

    # الشروط
    send_feedback = False
    if usage_count[uid] == 1:  # أول استخدام
        send_feedback = True
    elif usage_count[uid] % 5 == 0:  # كل 5 مرات
        # شرط مرور يوم كامل
        last_time = last_feedback_time.get(uid)
        if not last_time or (now - last_time) > timedelta(days=1):
            send_feedback = True

    if send_feedback:
        rating_markup = types.InlineKeyboardMarkup()
        rating_markup.row(
            types.InlineKeyboardButton("⭐ 1", callback_data="rate_1"),
            types.InlineKeyboardButton("⭐ 2", callback_data="rate_2"),
            types.InlineKeyboardButton("⭐ 3", callback_data="rate_3")
        )
        rating_markup.row(
            types.InlineKeyboardButton("⭐ 4", callback_data="rate_4"),
            types.InlineKeyboardButton("⭐ 5", callback_data="rate_5"),
            types.InlineKeyboardButton("تجاهل", callback_data="rate_ignore")
        )

        bot.send_message(
            chat_id,
            "✨ كيف كانت تجربتك مع TestGenie ؟\n\nاختر عدد النجوم للتقييم:",
            reply_markup=rating_markup
        )

        last_feedback_time[uid] = now


def send_daily_report():
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            today = datetime.date.today().isoformat()

            cursor.execute("""
                SELECT tests_generated, files_processed, new_users, channel_users, external_users
                FROM daily_stats
                WHERE date=?
            """, (today,))
            row = cursor.fetchone()
            if row:
                tests, files, new_users, channel_users, external_users = row
                msg = f"📊 تقرير اليوم ({today}):\n" \
                      f"📝 الاختبارات المولدة: {tests}\n" \
                      f"📂 الملفات المعالجة: {files}\n" \
                      f"👥 المستخدمون الجدد: {new_users} (قنوات: {channel_users} | خارجي: {external_users})"
                bot3.send_message(ADMIN_ID, msg)
    except Exception as e:
        logging.error(f"❌ Error sending daily report: {e}")


def send_top_users_report(top_n: int = 5):
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT user_id, tests_generated, files_processed FROM top_users ORDER BY tests_generated DESC LIMIT ?", (top_n,))
            rows = cursor.fetchall()
            msg = "🏆 أفضل المستخدمين:\n"
            for idx, (uid, tests, files) in enumerate(rows, 1):
                msg += f"{idx}. UserID: {uid} | Tests: {tests} | Files: {files}\n"
            bot3.send_message(ADMIN_ID, msg)
    except Exception as e:
        logging.error(f"❌ Error sending top users report: {e}")


# -------------------------------------------------------------------
# -------- معالجات الطلبات -------------------------------------------------------------------

# إعداد قائمة الانتظار والتحكم
request_queue = queue.Queue(maxsize=200)
semaphore = threading.Semaphore(5)
num_workers = 5







#!/usr/bin/env python3
# -*- coding: utf-8 -*-


# --- إعداد المفاتيح والعمل

session = {}  # <--- أضف هذا السطر
user_files = {}
# 1. إعداد Google Gemini
gemini_model = None
if GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel('gemini-1.5-flash')
        logging.info("✅ 1. Gemini configured successfully")
    except Exception as e:
        logging.warning(f"⚠️ Could not configure Gemini: {e}")

# 2. إعداد Groq
groq_client = None
if GROQ_API_KEY:
    try:
        groq_client = Groq(api_key=GROQ_API_KEY)
        logging.info("✅ 2. Groq configured successfully")
    except Exception as e:
        logging.warning(f"⚠️ Could not configure Groq: {e}")

# 3. إعداد OpenRouter (سيتم استخدامه لنموذجين مختلفين)
if OPENROUTER_API_KEY:
    logging.info("✅ 3. OpenRouter is ready")

# 4. إعداد Cohere
cohere_client = None
if COHERE_API_KEY:
    try:
        cohere_client = cohere.Client(COHERE_API_KEY)
        logging.info("✅ 4. Cohere configured successfully")
    except Exception as e:
        logging.warning(f"⚠️ Could not configure Cohere: {e}")


# --- الدالة الموحدة لتوليد الردود ---

def generate_gemini_response(prompt: str) -> str:
    """
    Tries to generate a response by attempting a chain of services silently.
    It logs errors for the developer but does not send progress messages to the user.
    """
    timeout_seconds = 45

    # 1️⃣ OpenRouter - Nous Hermes 2 (أفضل دعم للعربية)
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 1. OpenRouter (Nous Hermes 2)...")
            headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # ← غيّر هذا إلى رابط البوت
            "X-Title": "AI Quiz Bot"
            }
            model_identifier = "nousresearch/nous-hermes-2-mistral:free"
            response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
                json={
                "model": model_identifier,
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("✅ Success with OpenRouter (Nous Hermes 2).")
            return result_text
        except Exception as e:
            logging.warning(f"❌ OpenRouter (Nous Hermes 2) failed: {e}")

    # 2️⃣ Groq (LLaMA 3)
    if groq_client:
        try:
            logging.info("Attempting request with: 2. Groq (LLaMA 3)...")
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-8b-8192",
                temperature=0.7,
                timeout=timeout_seconds
            )
            if chat_completion.choices[0].message.content:
                logging.info("✅ Success with Groq.")
                return chat_completion.choices[0].message.content
            else:
                logging.warning("❌ Groq returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"❌ Groq failed: {e}")

    # 3️⃣ OpenRouter - Gemma
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 3. OpenRouter (Gemma)...")
            headers = {
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # Replace with your bot's link
                "X-Title": "AI Quiz Bot"
            }
            model_identifier = "google/gemma-7b-it:free"
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                json={"model": model_identifier, "messages": [{"role": "user", "content": prompt}]},
                timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("✅ Success with OpenRouter (Gemma).")
            return result_text
        except Exception as e:
            logging.warning(f"❌ OpenRouter (Gemma) failed: {e}")

    # 4️⃣ Google Gemini
    if gemini_model:
        try:
            logging.info("Attempting request with: 4. Google Gemini...")
            request_options = {"timeout": timeout_seconds}
            response = gemini_model.generate_content(prompt, request_options=request_options)
            if response.text:
                logging.info("✅ Success with Gemini.")
                return response.text
            else:
                logging.warning("❌ Gemini returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"❌ Gemini failed: {e}")

    # 5️⃣ Cohere
    if cohere_client:
        try:
            logging.info("Attempting request with: 5. Cohere...")
            response = cohere_client.chat(model='command-r', message=prompt)
            logging.info("✅ Success with Cohere.")
            return response.text
        except Exception as e:
            logging.warning(f"❌ Cohere failed: {e}")

    # 🚫 All models failed
    logging.error("❌ All API providers failed. Returning empty string.")
    return ""


def generate_smart_response(prompt: str) -> str:
    """
    Tries to generate a response by attempting a chain of services silently.
    It logs errors for the developer but does not send progress messages to the user.
    """
    timeout_seconds = 45


    #  1️⃣ Cohere
    if cohere_client:
        try:
            logging.info("Attempting request with: 5. Cohere...")
            response = cohere_client.chat(model='command-r', message=prompt, temperature=0.8)
            logging.info("✅ Success with Cohere.")
            return response.text
        except Exception as e:
            logging.warning(f"❌ Cohere failed: {e}")



    # 2️⃣ Google Gemini
    if gemini_model:
        try:
            logging.info("Attempting request with: 4. Google Gemini...")
            request_options = {"timeout": timeout_seconds}
            response = gemini_model.generate_content(prompt, request_options=request_options, temperature=0.8)
            if response.text:
                logging.info("✅ Success with Gemini.")
                return response.text
            else:
                logging.warning("❌ Gemini returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"❌ Gemini failed: {e}")


    #  3️⃣  Groq (LLaMA 3)
    if groq_client:
        try:
            logging.info("Attempting request with: 2. Groq (LLaMA 3)...")
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-8b-8192",
                temperature=0.8,
                timeout=timeout_seconds
            )
            if chat_completion.choices[0].message.content:
                logging.info("✅ Success with Groq.")
                return chat_completion.choices[0].message.content
            else:
                logging.warning("❌ Groq returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"❌ Groq failed: {e}")

    # 4️⃣# 5️⃣ OpenRouter - Gemma
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 3. OpenRouter (Gemma)...")
            headers = {
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # Replace with your bot's link
                "X-Title": "AI Quiz Bot"
            }
            model_identifier = "google/gemma-7b-it:free"
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                json={"model": model_identifier, "messages": [{"role": "user", "content": prompt}]},
                timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("✅ Success with OpenRouter (Gemma).")
            return result_text
        except Exception as e:
            logging.warning(f"❌ OpenRouter (Gemma) failed: {e}")

    # 🚫 All models failed
    logging.error("❌ All API providers failed. Returning empty string.")
    return ""

# -------------------------------------------------------------------
#                 OCR + language detection & translation 
# -------------------------------------------------------------------

def translate_text(text, source='en', target='ar'):
    url = 'https://libretranslate.de/translate'
    payload = {
        'q': text,
        'source': source,
        'target': target,
        'format': 'text'
    }
    try:
        response = requests.post(url, data=payload)
        return response.json()['translatedText']
    except Exception as e:
        print("ترجمة فشلت:", e)
        return text  # fallback


from flask import Flask, render_template, session, request, redirect, url_for

def save_user_major(user_id, major):
    with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO users (user_id, major)
            VALUES (?, ?)
            ON CONFLICT(user_id) DO UPDATE SET major=excluded.major
        """, (user_id, major))
        conn.commit()


from langdetect import detect

def detect_language(text: str) -> str:
    try:
        lang = detect(text)
        return lang
    except:
        return "unknown"

def detect_language_from_filename(filename: str) -> str:
    """
    يحاول تحديد اللغة المناسبة من اسم الملف.
    إذا احتوى على حروف عربية → يرجّح العربية.
    خلاف ذلك → الإنجليزية.
    """
    for char in filename:
        if '\u0600' <= char <= '\u06FF':  # نطاق الحروف العربية
            return "ara"
    return "eng"

import os
import fitz  # PyMuPDF
import logging
import requests
from pptx import Presentation
from tempfile import NamedTemporaryFile

# ---- OCR Space Integration ----
def extract_text_with_ocr_space(file_path: str, api_key="helloworld", language="eng") -> tuple:
    """
    Uses OCR.Space API to extract text from an image or scanned PDF.
    Returns: (text, debug_info)
    """
    url = 'https://api.ocr.space/parse/image'
    with open(file_path, 'rb') as f:
        response = requests.post(
            url,
            files={"file": f},
            data={
                "apikey": api_key,
                "language": language,
                "isOverlayRequired": False,
                "OCREngine": 2
            },
        )

    try:
        result = response.json()
        if result.get("IsErroredOnProcessing"):
            error_msg = result.get("ErrorMessage", "Unknown OCR error")
            return "", f"[OCR ERROR] {error_msg}"
        
        parsed = result.get("ParsedResults")
        if not parsed:
            return "", "[OCR ERROR] No ParsedResults returned."

        text = parsed[0].get("ParsedText", "").strip()
        return text, f"[OCR DEBUG] Length: {len(text)} | Excerpt: {text[:100]}"
    
    except Exception as e:
        return "", f"[OCR EXCEPTION] {e}"


# ---- PDF Split + OCR ----
def extract_text_from_pdf_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Splits a PDF into chunks of 3 pages (OCR.Space free limit),
    sends each chunk separately, and concatenates the extracted text.
    """
    try:
        doc = fitz.open(path)
        all_text = []
        # تقسيم كل 3 صفحات في ملف مؤقت
        for i in range(0, len(doc), 3):
            subdoc = fitz.open()  # ملف جديد مؤقت
            for j in range(i, min(i+3, len(doc))):
                subdoc.insert_pdf(doc, from_page=j, to_page=j)
            
            with NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                temp_path = tmp.name
                subdoc.save(temp_path)
                subdoc.close()
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)
            
            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PDF with OCR: {e}")
        return ""


# ---- PPTX Split + OCR ----
def extract_text_from_pptx_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Converts PPTX slides into smaller chunks (3 slides per file),
    sends each chunk separately to OCR.Space, and concatenates the text.
    """
    try:
        prs = Presentation(path)
        all_text = []

        # تقسيم العرض كل 3 شرائح
        for i in range(0, len(prs.slides), 3):
            new_ppt = Presentation()
            # إضافة تخطيط فارغ (مطلوب لعمل نسخ الشرائح)
            blank_layout = new_ppt.slide_layouts[6]

            for j in range(i, min(i+3, len(prs.slides))):
                slide = prs.slides[j]
                new_slide = new_ppt.slides.add_slide(blank_layout)
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        textbox = new_slide.shapes.add_textbox(left=0, top=0, width=new_ppt.slide_width, height=100)
                        textbox.text = shape.text
            
            with NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                temp_path = tmp.name
                new_ppt.save(temp_path)
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR PPTX chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)

            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PPTX with OCR: {e}")
        return ""
# -------------------------------------------------------------------
#                  Logging & Database Setup
# -------------------------------------------------------------------

logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s [%(levelname)s] %(message)s")


import sqlite3

def init_medical_db(db_path='medical_quizzes.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()

    cursor.execute('''
    CREATE TABLE IF NOT EXISTS quizzes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT NOT NULL,
        questions_json TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        user_id INTEGER,
        is_active BOOLEAN DEFAULT 1
    )
    ''')
    conn.commit()
    conn.close()


def init_user_quiz_db(db_path='quiz_users.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()

    
    # جدول المستخدمين
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
        user_id     INTEGER PRIMARY KEY,
        major       TEXT,
        native_lang TEXT DEFAULT 'ar',
        quiz_count  INTEGER DEFAULT 0,
        last_reset  TEXT
    )
    """)

    # جدول الأسئلة المقترحة من المستخدمين للعبة الاستنتاج
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS inference_questions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        question TEXT NOT NULL,
        options TEXT NOT NULL,         -- سيتم تخزينها كسلسلة JSON
        correct_index INTEGER NOT NULL,
        submitted_by INTEGER,
        approved INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS game_attempts (
        user_id INTEGER,
        game_type TEXT,
        date TEXT,
        PRIMARY KEY (user_id, game_type)
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS recent_questions (
        user_id INTEGER,
        game_type TEXT,
        question TEXT,
        added_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS user_quizzes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        quiz_data TEXT NOT NULL,
        quiz_code TEXT UNIQUE NOT NULL,
        created_at TEXT NOT NULL,
        is_active BOOLEAN DEFAULT 1
    )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS quiz_shares (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        quiz_code TEXT NOT NULL,
        shared_by_user_id INTEGER NOT NULL,
        shared_by_name TEXT,
        shared_at TEXT NOT NULL
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS sample_quizzes (
        quiz_code TEXT PRIMARY KEY,
        quiz_data TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)

    # مثال إضافة أعمدة جديدة بحذر
    # الأعمدة الجديدة
    new_columns = [
        ("score", "INTEGER"),
        ("total", "INTEGER"),
        ("timestamp", "TEXT"),
        ("owner_name", "TEXT")  # لتخزين اسم المالك
        ]
# التأكد من وجود العمود قبل إضافته
    for col_name, col_type in new_columns:
        cursor.execute(f"PRAGMA table_info(user_quizzes)")
        existing_cols = [row[1] for row in cursor.fetchall()]
        if col_name not in existing_cols:
            cursor.execute(f"ALTER TABLE user_quizzes ADD COLUMN {col_name} {col_type}")

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS bot_users (
        user_id INTEGER PRIMARY KEY,
        is_channel_user BOOLEAN DEFAULT 0, -- أعضاء القنوات
        is_external_user BOOLEAN DEFAULT 0 -- مستخدم جديد خارجي
    )
    """)
    
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS stat (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0,
        total_users INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS daily_stats (
        date TEXT PRIMARY KEY,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0,
        new_users INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS top_users (
        user_id INTEGER PRIMARY KEY,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0
    )
    """)
    conn.commit()
    conn.close()

    ALTER TABLE daily_stats ADD COLUMN channel_users INTEGER DEFAULT 0;
    ALTER TABLE daily_stats ADD COLUMN external_users INTEGER DEFAULT 0;
    
def init_request_db(db_path='requests.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cur = conn.cursor()
    cur.execute('''
        CREATE TABLE IF NOT EXISTS requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            username TEXT,
            file_id TEXT,
            status TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()

  
import sqlite3, logging

def fetch_user_major(uid, db_path="quiz_users.db"):
    try:
        with sqlite3.connect(db_path, check_same_thread=False) as conn:
            cur = conn.cursor()
            cur.execute("SELECT major FROM users WHERE user_id=?", (uid,))
            row = cur.fetchone()
        return row[0] if row else "General"
    except Exception:
        logging.exception("fetch_user_major failed")
        return "General"



def init_all_dbs():
    init_medical_db()
    init_user_quiz_db()
    init_request_db()


# ---------------------------
# ---- get basic statics -----
# -------------    ---------

import sqlite3

DB_NAME = "quiz_users.db"

def get_tests_generated():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT tests_generated FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


def get_files_processed():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT files_processed FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


def get_total_users():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT total_users FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


# ----------------------------
# ----    flask config  ----------------------------

# واجهة Flask للفحص
app = Flask(__name__)

@app.route('/')
def home():
    return render_template('index.html')
    

@app.route('/quiz/<int:quiz_id>')
def show_quiz(quiz_id):
    conn = sqlite3.connect('medical_quizzes.db')
    cursor = conn.cursor()
    
    cursor.execute('SELECT title, questions_json FROM quizzes WHERE id = ?', (quiz_id,))
    quiz = cursor.fetchone()
    conn.close()
    
    if quiz:
        quiz_data = {
            'title': quiz[0],
            'questions': json.loads(quiz[1])
        }
        return render_template('quiz.html', quiz=quiz_data)
    else:
        return "الاختبار غير موجود", 404




@app.route('/supportme')
def supportme():
    # هنا تقوم بجلب البيانات الحقيقية من قاعدة البيانات أو أي مصدر آخر
    # هذه مجرد قيم افتراضية
    tests_generated = get_tests_generated()
    files_processed = get_files_processed()
    total_users = get_total_users()

    return render_template(
        'supportme.html',
        tests_generated=tests_generated,
        files_processed=files_processed,
        total_users=total_users
    )








# track temporary state for custom-major input
user_states = {}
usage_count = {}
import sqlite3, logging




# quene

def save_request(msg, db_path='requests.db'):
    try:
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        file_id = getattr(getattr(msg, "document", None), "file_id", None)
        cur.execute(
            'INSERT INTO requests (user_id, username, file_id, status) VALUES (?, ?, ?, ?)',
            (msg.from_user.id, msg.from_user.username or "", file_id, 'pending')
        )
        conn.commit()
    except Exception:
        logging.exception("save_request failed")
    finally:
        try: conn.close()
        except: pass

def update_request_status(file_id, new_status, db_path='requests.db'):
    try:
        if not file_id:
            return
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        cur.execute('UPDATE requests SET status=? WHERE file_id=?', (new_status, file_id))
        conn.commit()
    except Exception:
        logging.exception("update_request_status failed")
    finally:
        try: conn.close()
        except: pass



# دالة العامل
def worker():
    while True:
        item = request_queue.get()
        try:
            if isinstance(item, tuple) and len(item) == 2:
                msg, sent_msg = item
                process_message(msg, message_id=sent_msg.message_id, chat_id=sent_msg.chat.id)
            else:
                # عنصر واحد: فقط الرسالة
                msg = item
                process_message(msg)
        except Exception:
            logging.exception("[WORKER ERROR]")
        finally:
            request_queue.task_done()
# تشغيل العمال في الخلفية
def start_workers():
    for _ in range(num_workers):
        threading.Thread(target=worker, daemon=True).start()
    logging.info("Workers started: %s", num_workers)




def safe_edit_or_send(text, chat_id, message_id, parse_mode="HTML"):
    try:
        if chat_id and message_id:
            return bot.edit_message_text(
                text, chat_id=chat_id, message_id=message_id, parse_mode=parse_mode
            )
    except Exception as e:
        logging.warning("edit_message_text failed (%s), fallback to send_message", e)
    return bot.send_message(chat_id, text, parse_mode=parse_mode)



def is_request_already_queued(file_id=None, user_id=None, message_id=None, db_path='requests.db'):
    try:
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        if file_id:
            cur.execute("SELECT status FROM requests WHERE file_id=? ORDER BY id DESC LIMIT 1", (file_id,))
        else:
            cur.execute("SELECT status FROM requests WHERE user_id=? AND message_id=? ORDER BY id DESC LIMIT 1", (user_id, message_id))
        row = cur.fetchone()
        conn.close()
        if row and row[0] in ('pending','processing'):
            return True
        return False
    except Exception:
        logging.exception("is_request_already_queued failed")
        return False


# ------------------------------------
# ------------------------------------------------- Stat Management----
# ----------------------------------------------------------------------

def add_external_user(uid: int):
    with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (uid,))
        if cursor.fetchone()[0] == 0:
            cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (uid, 1))
            conn.commit()


def update_files_and_users(uid: int = None, files_count: int = 1):
    """
    - files_count: عدد الملفات الجديدة
    - uid: user_id للمستخدم الجديد الخارجي (اختياري)
    """
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # 1) تحديث المستخدم الجديد إذا تم تمرير uid
            if uid:
                cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (uid,))
                if cursor.fetchone()[0] == 0:
                    # إضافة المستخدم الخارجي الجديد
                    cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (uid, 1))

            # 2) حساب عدد مستخدمي القنوات
            env_channels = os.getenv("ALLOWED_CHANNELS", "")
            channel_users_count = 0
            if env_channels.strip():
                allowed_channels = set(map(int, env_channels.split(",")))
                for channel_id in allowed_channels:
                    try:
                        chat = bot.get_chat(channel_id)
                        channel_users_count += bot.get_chat_members_count(channel_id)
                    except Exception as e:
                        logging.warning(f"⚠️ Failed to fetch channel members for {channel_id}: {e}")

            # 3) حساب عدد المستخدمين الخارجيين المسجلين في الجدول
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_external_user=1")
            external_users_count = cursor.fetchone()[0]

            total_users = channel_users_count + external_users_count

            # 4) تحديث الإحصائيات العامة في جدول stat
            cursor.execute("""
                UPDATE stat
                SET files_processed = files_processed + ?,
                    total_users = ?
                WHERE id = 1
            """, (files_count, total_users))

            conn.commit()
    except Exception as e:
        logging.error(f"❌ Error updating files and users: {e}")
        

def update_daily_stats(date_str: str = None, tests: int = 0, files: int = 0):
    """
    تحديث الإحصائيات اليومية:
    - يحسب المستخدمين الجدد تلقائيًا من جدول users
    """
    import datetime
    date_str = date_str or datetime.date.today().isoformat()
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # حساب المستخدمين الجدد حسب النوع
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_channel_user=1")
            channel_users = cursor.fetchone()[0]

            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_external_user=1")
            external_users = cursor.fetchone()[0]

            total_new_users = channel_users + external_users

            # إدراج أو تحديث سجل اليوم
            cursor.execute("SELECT COUNT(*) FROM daily_stats WHERE date=?", (date_str,))
            if cursor.fetchone()[0] == 0:
                cursor.execute("""
                    INSERT INTO daily_stats (date, tests_generated, files_processed, new_users, channel_users, external_users)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (date_str, tests, files, total_new_users, channel_users, external_users))
            else:
                cursor.execute("""
                    UPDATE daily_stats
                    SET tests_generated = tests_generated + ?,
                        files_processed = files_processed + ?,
                        new_users = ?,
                        channel_users = ?,
                        external_users = ?
                    WHERE date = ?
                """, (tests, files, total_new_users, channel_users, external_users, date_str))

            conn.commit()
    except Exception as e:
        logging.error(f"❌ Error updating daily stats: {e}")
    

def update_top_user(user_id: int, tests: int = 0, files: int = 0):
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # تأكد أن المستخدم موجود في جدول users
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (user_id,))
            if cursor.fetchone()[0] == 0:
                # إضافة المستخدم تلقائيًا كمستخدم خارجي
                cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (user_id, 1))

            # تحديث جدول top_users
            cursor.execute("SELECT COUNT(*) FROM top_users WHERE user_id=?", (user_id,))
            if cursor.fetchone()[0] == 0:
                cursor.execute("INSERT INTO top_users (user_id, tests_generated, files_processed) VALUES (?, ?, ?)",
                               (user_id, tests, files))
            else:
                cursor.execute("""
                    UPDATE top_users
                    SET tests_generated = tests_generated + ?,
                        files_processed = files_processed + ?
                    WHERE user_id = ?
                """, (tests, files, user_id))

            conn.commit()
    except Exception as e:
        logging.error(f"❌ Error updating top user: {e}")
# -------------------------------------------------------------------
#                     Text Extraction & OCR
# -------------------------------------------------------------------

def extract_text_from_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        text = "\n".join([page.get_text() for page in doc])
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting PDF text: {e}")
        return ""
    # fallback to PyMuPDF text extraction
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])
# أضف هذه الدالة في قسم Text Extraction & OCR
def extract_text_from_docx(path: str) -> str:
    try:
        doc = docx.Document(path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting DOCX text: {e}")
        return ""

# ويجب أيضاً تعريف دالة لملفات txt
def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logging.error(f"Error extracting TXT text: {e}")
        return ""
        
def is_text_empty(text: str) -> bool:
    return not text or len(text.strip()) < 30  # يمكن تعديل الحد حسب تجربتك


def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text).strip()
    except Exception as e:
        logging.error(f"Error extracting PPTX text: {e}")
        return ""

def split_text(content, chunk_size=3500):
    return [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]

def summarize_long_text(content: str) -> str:
    """
    Summarize the content in its original language (Arabic or English) using educational style.
    """
    lang = detect_language(content[:1000])  # نكتفي بأول 1000 حرف للتحليل
    print(f"[DEBUG] Detected language: {lang}")

    if lang.startswith("ar"):
        summary_prompt = (
            "أنت مساعد تعليمي محترف. قم بتلخيص المحتوى التالي بأسلوب تعليمي منظم وواضح باللغة العربية فقط،"
            " مع الحفاظ على النقاط المفيدة والتفاصيل التي يمكن استخدامها لاحقًا لصنع أسئلة أو بطاقات تعليمية."
            "\n\nالمحتوى:\n{chunk}"
        )
        merge_prompt = (
            "فيما يلي مجموعة من الملخصات الجزئية لمحتوى تعليمي. قم بدمجها في ملخص نهائي شامل ومترابط وواضح"
            " باللغة العربية، مع الحفاظ على التفاصيل المفيدة التي تساعد على فهم المحتوى أو إعداد اختبار منه.\n\n{merged}"
        )
    else:
        summary_prompt = (
            "You are a professional educational assistant. Summarize the following content into a clear and concise educational explanation in **English only**.\n\n"
            "- Preserve factual details and key definitions.\n"
            "- Avoid vague sentences or repetition.\n"
            "- Keep the original language (do not translate).\n\nContent:\n{chunk}"
        )
        merge_prompt = (
            "You are an educational summarizer. Merge the following partial summaries into one final, well-structured summary in **English**, preserving all useful learning content.\n\n{merged}"
        )

    chunks = split_text(content)
    partial_summaries = []

    for i, chunk in enumerate(chunks):
        prompt = summary_prompt.format(chunk=chunk)
        summary = generate_smart_response(prompt.strip())
        partial_summaries.append(summary)

    merged_summary = "\n".join(partial_summaries)
    final_prompt = merge_prompt.format(merged=merged_summary)
    return generate_smart_response(final_prompt.strip())
    

def parse_ai_json(raw_text: str) -> dict | None:

    # 1. فكُّ هاربات Unicode (\u0627 → ا)
    def _unescape(match):
        code = match.group(1)
        return chr(int(code, 16))
    text = re.sub(r'\\u([0-9A-Fa-f]{4})', _unescape, raw_text)

    # 2. اجتزء أول كتلة JSON (من { إلى })
    m = re.search(r'\{[\s\S]*\}', text)
    json_text = m.group(0) if m else text

    # 3. حاول التحميل
    for attempt in (json_text, text):
        try:
            data = json.loads(attempt)
            break
        except json.JSONDecodeError:
            data = None
    if not data:
        return None

    # 4. التحقق من بنية الـ dict
    if not all(k in data for k in ("question", "options", "correct_index")):
        return None

    # 5. التأكد من أن options قائمة وصالحة
    if not isinstance(data["options"], list) or len(data["options"]) < 2:
        return None

    # 6. التأكد من correct_index
    ci = data["correct_index"]
    if not isinstance(ci, int) or ci < 0 or ci >= len(data["options"]):
        return None

    return data

def generate_game(prompt, user_id=0, translate_all=False, translate_question=False):
    if user_id == ADMIN_ID or can_generate(user_id):  # <-- التحقق هنا
        raw_response = generate_smart_response(prompt)
    else:
        raw_response = generate_gemini_response(prompt)
        
    game_data = parse_ai_json(raw_response)

    if not game_data:
        raise ValueError("فشل استخراج بيانات اللعبة")

    if translate_all:
        # ترجمة السؤال
        if 'question' in game_data:
            game_data['question'] = translate_text(game_data['question'], source='en', target='ar')

        # ترجمة كل الخيارات
        if 'options' in game_data and isinstance(game_data['options'], list):
            game_data['options'] = [
                translate_text(option, source='en', target='ar') for option in game_data['options']
            ]

    elif translate_question:
        # ترجمة السؤال فقط
        if 'question' in game_data:
            game_data['question'] = translate_text(game_data['question'], source='en', target='ar')

    return game_data

import genanki
import time
import uuid

def save_cards_to_apkg(cards: List[Dict], filename: str = 'anki_flashcards.apkg', deck_name: str = "My Flashcards"):
    model = genanki.Model(
        1607392319,
        'Simple Model with Tags',
        fields=[
            {'name': 'Front'},
            {'name': 'Back'},
            {'name': 'Tag'}
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Front}}<br><small style="color:gray">{{Tag}}</small>',
                'afmt': '{{FrontSide}}<hr id="answer">{{Back}}',
            },
        ]
    )

    deck = genanki.Deck(
        deck_id=int(str(uuid.uuid4().int)[:9]),
        name=deck_name
    )

    temp_dir = tempfile.mkdtemp()  # إنشاء مجلد مؤقت خاص بهذه الدفعة
    media_files = []

    for idx, card in enumerate(cards, start=1):
        try:
            front = card.get('front', '').strip()
            back = card.get('back', '').strip()
            tag = card.get('tag', '').strip()
            image_hint = card.get('image_hint', '').strip()

            if not front or not back:
                continue

            # البحث عن الصورة وتحميلها إذا كان هناك تلميح
            if image_hint:
                image_filename = search_and_download_image(image_hint, temp_dir)
                if image_filename:
                    media_files.append(os.path.join(temp_dir, image_filename))
                    back += f"<br><img src='{image_filename}' style='max-height:220px;'>"

            # إنشاء البطاقة
            note = genanki.Note(model=model, fields=[front, back, tag])
            deck.add_note(note)

        except Exception as e:
            logging.error(f"❌ خطأ في البطاقة #{idx}: {e}")
            continue

    # إنشاء الحزمة النهائية
    try:
        package = genanki.Package(deck)
        if media_files:
            package.media_files = media_files
        package.write_to_file(filename)
        logging.info(f"✅ تم إنشاء ملف Anki: {filename}")
        return filename
    except Exception as e:
        logging.error(f"❌ خطأ في إنشاء ملف Anki: {e}")
        return None
    finally:
        # تنظيف الملفات المؤقتة (اختياري)
        try:
            shutil.rmtree(temp_dir)
        except:
            pass
        


def parse_manual_anki_input(text):
    cards = []
    lines = [line.strip() for line in text.strip().split('\n')]
    current_card = []

    for line in lines:
        if line == "":
            if current_card:
                if len(current_card) >= 2:
                    card = {
                        "front": current_card[0],
                        "back": current_card[1],
                        "tag": current_card[2] if len(current_card) > 2 else ""
                    }
                    cards.append(card)
                current_card = []
        else:
            current_card.append(line)

    if current_card and len(current_card) >= 2:
        card = {
            "front": current_card[0],
            "back": current_card[1],
            "tag": current_card[2] if len(current_card) > 2 else ""
        }
        cards.append(card)

    return cards
    
# -------------------------------------------------------------------
#                     Quota Management
# -------------------------------------------------------------------
def add_recent_question(user_id, game_type, question):
    with sqlite3.connect("quiz_users.db") as conn:
        cursor = conn.cursor()
        
        # إدخال السؤال الجديد
        cursor.execute("""
        INSERT INTO recent_questions (user_id, game_type, question) 
        VALUES (?, ?, ?)
        """, (user_id, game_type, question))
        
        # حذف الأقدم إذا تجاوز 10 أسئلة
        cursor.execute("""
        DELETE FROM recent_questions
        WHERE user_id = ? AND game_type = ?
        AND question NOT IN (
            SELECT question FROM recent_questions
            WHERE user_id = ? AND game_type = ?
            ORDER BY added_at DESC
            LIMIT 10
        )
        """, (user_id, game_type, user_id, game_type))

        conn.commit()

def get_recent_questions(user_id, game_type):
    with sqlite3.connect("quiz_users.db") as conn:
        cursor = conn.cursor()
        cursor.execute("""
        SELECT question FROM recent_questions
        WHERE user_id = ? AND game_type = ?
        ORDER BY added_at DESC
        LIMIT 10
        """, (user_id, game_type))
        rows = cursor.fetchall()
        return [row[0] for row in rows]


def reset_if_needed(user_id: int):
    this_month = datetime.now().strftime("%Y-%m")
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            
            cursor.execute("SELECT last_reset FROM users WHERE user_id = ?", (user_id,))
            row = cursor.fetchone()
            
            if not row or row[0] != this_month:
                cursor.execute("""
                    INSERT OR REPLACE INTO users(user_id, major, quiz_count, last_reset)
                    VALUES (?, COALESCE((SELECT major FROM users WHERE user_id=?), ''), 0, ?)
                """, (user_id, user_id, this_month))
                conn.commit()
    except Exception as e:
        logging.error(f"🚫 خطأ في reset_if_needed للمستخدم {user_id}: {e}")






MAX_FREE_ATTEMPTS = 3  # 👈 عدلها حسب ما تريد
def can_generate(user_id: int) -> bool:
    # السماح للأدمن دائماً
    if user_id == ADMIN_ID:
        return True
    
    # استخدام 'with' يضمن فتح وإغلاق الاتصال بشكل آمن
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            
            today = datetime.utcnow().date()
            cursor.execute("SELECT quiz_count, last_reset FROM users WHERE user_id = ?", (user_id,))
            row = cursor.fetchone()
            
            quiz_count = 0  # قيمة افتراضية

            if not row:
                # مستخدم جديد
                cursor.execute("INSERT INTO users (user_id, quiz_count, last_reset) VALUES (?, ?, ?)", (user_id, 0, today.isoformat()))
                conn.commit()
                quiz_count = 0
            else:
                quiz_count, last_reset = row
                last_reset_date = None
                try:
                    if last_reset:
                        last_reset_date = datetime.fromisoformat(last_reset).date()
                except Exception as e:
                    logging.warning(f"⚠️ last_reset غير صالحة للمستخدم {user_id}: {e}")
                    last_reset_date = None
                
                # تحقق من الدخول في شهر جديد
                if not last_reset_date or last_reset_date.month != today.month or last_reset_date.year != today.year:
                    cursor.execute("UPDATE users SET quiz_count = 0, last_reset = ? WHERE user_id = ?", (today.isoformat(), user_id))
                    conn.commit()
                    quiz_count = 0
            
            # التحقق من عدد المحاولات
            if quiz_count >= MAX_FREE_ATTEMPTS:
                # التحقق من القنوات المسموح بها إذا تجاوز الحد
                try:
                    raw = os.getenv("ALLOWED_CHANNELS", "")
                    if not raw.strip(): return False # إذا لم تكن هناك قنوات، لا تسمح
                    
                    allowed_channels = set(int(cid) for cid in raw.split(",") if cid.strip())
                    for channel_id in allowed_channels:
                        try:
                            member = bot.get_chat_member(chat_id=channel_id, user_id=user_id)
                            if member.status in ['member', 'administrator', 'creator']:
                                return True  # مستخدم مميز
                        except Exception as e:
                            logging.warning(f"⚠️ فشل التحقق من القناة {channel_id} للمستخدم {user_id}: {e}")
                    return False  # ليس عضوًا في أي قناة مسموح بها
                except Exception as e:
                    logging.error(f"🚫 خطأ في قراءة القنوات المسموح بها: {e}")
                    return False
            else:
                return True  # ضمن الحد المسموح به
    
    except Exception as e:
        logging.error(f"🚫 خطأ فادح في دالة can_generate: {e}")
        return False
        


from datetime import datetime


def increment_count(user_id: int):
    # لا تقم بزيادة العداد إذا كان المستخدم هو الأدمن أو مستخدم مميز
    if user_id == ADMIN_ID:
        bot.send_message(ADMIN_ID, "✨ (وضع الأدمن: لم يتم احتساب هذه المحاولة)")
        return
    
    try:
        # التحقق من القنوات المميزة
        raw = os.getenv("ALLOWED_CHANNELS", "")
        allowed_channels = set(int(cid) for cid in raw.split(",") if cid.strip())
        for channel_id in allowed_channels:
            try:
                member = bot.get_chat_member(chat_id=channel_id, user_id=user_id)
                if member.status in ['member', 'administrator', 'creator']:
                    return  # مستخدم مميز
            except Exception as e:
                logging.warning(f"⚠️ فشل التحقق من القناة {channel_id} للمستخدم {user_id}: {e}")
        
        # إذا لم يكن مميز → زيادة العداد
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            cursor.execute("UPDATE users SET quiz_count = quiz_count + 1 WHERE user_id = ?", (user_id,))
            conn.commit()
    
    except Exception as e:
        logging.error(f"🚫 خطأ في increment_count للمستخدم {user_id}: {e}")

def can_play_game_today(user_id: int, game_type: str) -> bool:
    if str(user_id) == str(ADMIN_ID):  # مقارنة آمنة لأن ADMIN_ID أحيانًا يكون str
        return True

    today = str(date.today())
    cursor.execute(
        "SELECT 1 FROM game_attempts WHERE user_id = ? AND game_type = ? AND date = ?",
        (user_id, game_type, today)
    )
    return cursor.fetchone() is None

def record_game_attempt(user_id: int, game_type: str):
    if str(user_id) == str(ADMIN_ID):
        return  # لا تسجل للأدمن

    today = str(date.today())
    cursor.execute(
        "INSERT OR REPLACE INTO game_attempts(user_id, game_type, date) VALUES (?, ?, ?)",
        (user_id, game_type, today)
    )
    conn.commit()
from collections import defaultdict

# تخزين مؤقت في الذاكرة
game_states = defaultdict(dict)  # {user_id: {game_type: count}}

def get_question_count(user_id, game_type):
    return game_states.get(user_id, {}).get(game_type, 0)

def increment_question_count(user_id, game_type):
    game_states[user_id][game_type] = game_states.get(user_id, {}).get(game_type, 0) + 1
    
# -------------------------------------------------------------------
#                 Quiz Generation & Formatting
# -------------------------------------------------------------------

def extract_json_from_string(text: str) -> str:
    """
    Extracts a JSON string from a text that might contain markdown code blocks or other text.
    """
    # البحث عن بلوك JSON داخل ```json ... ```
    match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
    if match:
        return match.group(1).strip()

    # إذا لم يجد بلوك، ابحث عن أول '{' أو '[' وآخر '}' أو ']'
    start = -1
    end = -1
    
    # البحث عن بداية القائمة أو الكائن
    first_brace = text.find('{')
    first_bracket = text.find('[')
    
    if first_brace == -1:
        start = first_bracket
    elif first_bracket == -1:
        start = first_brace
    else:
        start = min(first_brace, first_bracket)

    # إذا لم يتم العثور على بداية، أرجع النص الأصلي
    if start == -1:
        return text

    # البحث عن نهاية القائمة أو الكائن
    last_brace = text.rfind('}')
    last_bracket = text.rfind(']')
    end = max(last_brace, last_bracket)

    # إذا تم العثور على بداية ونهاية، أرجع ما بينهما
    if end > start:
        return text[start:end+1].strip()
        
    # كخيار أخير، أرجع النص كما هو
    return text
    
def generate_quizzes_from_text(content: str, major: str, user_id: int, num_quizzes: int = 10):
    prompt = (
        f"You are a strict AI quiz generator. Your only task is to generate a JSON array of {num_quizzes} quiz questions "
        f"that are based **strictly and only** on the information explicitly stated in the following content.\n\n"
        "❗️Important Rules:\n"
        "- DO NOT invent, infer, or assume any information not clearly mentioned in the text.\n"
        "- If a concept is not explained or mentioned clearly in the content, DO NOT create a question about it.\n"
        "- Stay fully inside the boundaries of the content.\n"
        "- Every question must test **recall** or **recognition** from the provided text only, not general knowledge.\n"
        "- Questions must be varied: some fill-in-the-blank, some multiple-choice.\n"
        "- Include at most one True/False question.\n"
        "- All questions and answers must be in the same language as the content.\n"
        "- if the content language is arabic give the questions and answers in arabic.\n\n"
        "Each question must be an object with:\n"
        "- 'question': the question string\n"
        "- 'options': a list of exactly 4 answer options\n"
        "- 'correct_index': the index (0-3) of the correct answer in the options list\n"
        "- 'explanation': short sentence to explain **why this is the correct answer**, max 2 lines\n\n"
        "⚠️ Format Instructions:\n"
        "- ONLY return a raw JSON array. No markdown, no explanation, no formatting.\n"
        "- Do not include any introductory or closing text.\n"
        "- Ensure the JSON is valid and parsable.\n\n"
        f"Content:\n{content}"
    )

    # تحديد الدالة بناءً على صلاحية المستخدم
    if user_id == ADMIN_ID or can_generate(user_id):  # <-- التحقق هنا
        raw_response = generate_smart_response(prompt)
    else:
        raw_response = generate_gemini_response(prompt)
    
    # --- التعديل يبدأ هنا ---
    # 1. تنظيف الاستجابة لاستخراج الـ JSON
    clean_json_str = extract_json_from_string(raw_response)
    
    # 2. التحقق مما إذا كانت الاستجابة فارغة بعد التنظيف
    if not clean_json_str:
        logging.error(f"❌ JSON extraction failed. Raw output was:\n{raw_response}")
        return [] # أرجع قائمة فارغة بدلاً من رسالة خطأ

    try:
        # 3. محاولة تحليل السلسلة النظيفة
        quizzes_json = json.loads(clean_json_str)
        quizzes = []

        for item in quizzes_json:
            q = item.get("question", "").strip()
            opts = item.get("options", [])
            corr = item.get("correct_index", -1)
            expl = item.get("explanation", "").strip()

            if isinstance(q, str) and q and isinstance(opts, list) and len(opts) == 4 and isinstance(corr, int) and 0 <= corr < 4:
                quizzes.append((q, [str(opt).strip() for opt in opts], corr, expl))
            else:
                logging.warning(f"❌ Skipping invalid question structure: {item}")

        return quizzes

    except json.JSONDecodeError as e:
        logging.error(f"❌ JSON parsing failed: {e}\nCleaned string was:\n{clean_json_str}\nRaw output was:\n{raw_response}")
        return [] # أرجع قائمة فارغة عند الفشل
    # --- التعديل ينتهي هنا ---


def save_quiz_to_db(quiz_data, user_id):
    conn = sqlite3.connect('medical_quizzes.db')
    cursor = conn.cursor()
    
    cursor.execute('''
    INSERT INTO quizzes (title, questions_json, user_id)
    VALUES (?, ?, ?)
    ''', (quiz_data['title'], json.dumps(quiz_data['questions']), user_id))
    
    quiz_id = cursor.lastrowid
    conn.commit()
    conn.close()
    return quiz_id



import json
import logging
import requests

WIKIMEDIA_API = "https://commons.wikimedia.org/w/api.php"

def search_image_on_wikimedia(query: str) -> str:
    """
    البحث عن صورة من Wikimedia بناءً على الوصف
    """
    params = {
        "action": "query",
        "format": "json",
        "prop": "imageinfo",
        "generator": "search",
        "gsrsearch": query + " filetype:bitmap OR filetype:jpeg OR filetype:png",
        "gsrlimit": 1,
        "iiprop": "url",
    }
    try:
        r = requests.get(WIKIMEDIA_API, params=params)
        r.raise_for_status()
        data = r.json()

        if "query" in data and "pages" in data["query"]:
            page = next(iter(data["query"]["pages"].values()))
            if "imageinfo" in page:
                return page["imageinfo"][0]["url"]
        return ""
    except Exception as e:
        logging.error(f"❌ فشل البحث عن الصورة: {e}")
        return ""

def generate_Medical_quizzes(content: str, major: str, user_id: int, num_quizzes: int = 10):
    # (البرومبت المحسن من الخطوة 2 يجب وضعه هنا)
    prompt = (
        f"You are a medical education expert. Your task is to create a JSON-formatted quiz for {major} "
        "medical students (Year 3-4) based ONLY on the provided reference text.\n\n"
        "## EXTREMELY STRICT RULES:\n"
        f"1. You MUST STRICTLY generate {num_quizzes} questions. No more, no less.\n"
        "2. 90% multiple-choice (basic sciences), 10% problem-solving (an MCQ clinical case).\n"
        "3. Use only information from the reference text.\n"
        "4. Clinical questions must have realistic short scenarios (2-3 sentences).\n"
        "5. For any question that would benefit from an image, add an 'image_prompt' field.\n"
        "6. Language: English.\n"
        "7. CRITICAL: The 'questions' array MUST NOT be empty. If you cannot generate questions from the text, return an empty JSON object {} and nothing else.\n\n"
        "8. only create on clinical case in the form of MCQ and all questions should be Multiple choice questions"
        "## JSON OUTPUT STRUCTURE:\n"
        "{\n"
        "  \"title\": \"Medical Quiz in [major]\",\n"
        "  \"questions\": [\n"
        "    {\n"
        "      \"id\": 1,\n"
        "      \"type\": \"multiple_choice\",\n"
        "      \"question\": \"...\",\n"
        "      \"options\": [\"...\", \"...\", \"...\", \"...\"],\n"
        "      \"correct_index\": 0,\n"
        "      \"image_prompt\": \"Description of the image (optional)\"\n"
        "    }\n"
        "  ]\n"
        "}\n\n"
        f"## Reference text:\n{content}"
    )

    # --- آلية إعادة المحاولة (الخطوة 3) ---
    for attempt in range(3):  # سيحاول حتى 3 مرات
        logging.info(f"Attempt {attempt + 1} to generate medical quiz...")
        
        # اختيار النموذج
        if user_id == ADMIN_ID or can_generate(user_id):
            raw_response = generate_smart_response(prompt)
        else:
            raw_response = generate_gemini_response(prompt)
        
        clean_json_str = extract_json_from_string(raw_response)
        if not clean_json_str or clean_json_str == "{}":
            logging.warning(f"Attempt {attempt + 1} failed: AI returned empty response.")
            continue # انتقل للمحاولة التالية

        try:
            quiz_data = json.loads(clean_json_str)

            # --- التحقق المشدد (الخطوة 1) ---
            if "title" not in quiz_data or not quiz_data.get("questions"):
                raise ValueError("Invalid JSON: missing title or questions array is empty.")

            # --- المعالجة الناجحة ---
            for i, q in enumerate(quiz_data["questions"]):
                if "image_prompt" in q and q["image_prompt"].strip():
                    image_url = search_image_on_wikimedia(q["image_prompt"])
                    q["image_url"] = image_url if image_url else ""
                else:
                    q["image_url"] = ""
                if "id" not in q:
                    q["id"] = i + 1

            if major not in quiz_data["title"]:
                quiz_data["title"] = f"Medical Quiz in {major}"

            quiz_id = save_quiz_to_db(quiz_data, user_id)
            quiz_data["db_id"] = quiz_id

            logging.info(f"Successfully generated quiz on attempt {attempt + 1}.")
            return quiz_data # إرجاع البيانات الناجحة وإنهاء الدالة

        except (json.JSONDecodeError, ValueError) as e:
            # --- تسجيل الخطأ للتشخيص (الخطوة 4) ---
            logging.error(f"Attempt {attempt + 1} failed during JSON processing: {e}")
            logging.error(f"RAW AI RESPONSE WAS:\n{raw_response}\n")
            continue # انتقل للمحاولة التالية

    # إذا فشلت كل المحاولات
    logging.error("All 3 attempts to generate medical quiz failed.")
    return None
    


def generate_anki_cards_from_text(content: str, major: str = "General", user_id: int = 0, num_cards: int = 15) -> tuple:
    for attempt in range(3):  # تجربة حتى 3 مرات
        prompt = f"""
You are an AI assistant specialized in creating study flashcards.

🎯 Task:
Extract the most important {num_cards} points from the following content, and convert each into an **Anki-style flashcard**.

🔹 Rules:
- Each flashcard must include:
  - "front": a short question or hint.
  - "back": the detailed answer or explanation.
  - "tag": (optional) topic label like Grammar, Biology, Logic, etc.
- The front must be phrased to encourage recall (e.g. "What is...", "Define...", "How does...").
- Don't use Markdown, just clean plain text.
- Keep the cards diverse and helpful.
- Output must be a valid JSON **object** with two keys: "title" and "cards".

🚫 Important:
- Do NOT generate multiple choice or true/false questions.
- Only generate flashcards suitable for Anki with a front and a back.
- The flashcards must be written in the same language as the input content. If the content is in Arabic, answer in Arabic. If English, answer in English.

📘 Content to process (field: {major}):
{content}

✅ Example output format:
{{
  "title": "Basics of Organic Chemistry",
  "cards": [
    {{
      "front": "What is the function of mitochondria?",
      "back": "It is the powerhouse of the cell.",
      "tag": "Biology"
    }},
    {{
      "front": "ما هي الاستعارة؟",
      "back": "الاستعارة هي استخدام الكلمة في غير معناها الحقيقي لعلاقة مع قرينة مانعة.",
      "tag": "Literature"
    }}
  ]
}}
"""
        if user_id == ADMIN_ID or can_generate(user_id):  # <-- التحقق هنا
            raw_output = generate_smart_response(prompt)
        else:
            raw_output = generate_gemini_response(prompt)
            
        clean_json = extract_json_from_string(raw_output)

        try:
            data = json.loads(clean_json)
            title = data.get("title", "بطاقات تعليمية")
            card_list = data.get("cards", [])

            cards = []
            for item in card_list:
                front = item.get("front") or item.get("question")
                back = item.get("back") or item.get("answer")

                if isinstance(front, str) and isinstance(back, str) and front.strip() and back.strip():
                    cards.append({"front": front.strip(), "back": back.strip()})
                else:
                    logging.warning(f"❌ Skipping invalid card: {item}")

            if len(cards) >= 5:
                return cards, title

        except json.JSONDecodeError as e:
            logging.error(f"❌ Failed to parse Anki cards: {e}\nClean JSON:\n{clean_json}\nRaw:\n{raw_output}")

    return [], "بطاقات تعليمية"   



import json
import logging
import requests
from pptx import Presentation

WIKIMEDIA_API = "https://commons.wikimedia.org/w/api.php"
UNSPLASH_API = "https://api.unsplash.com/search/photos"
PEXELS_API = "https://api.pexels.com/v1/search"

# --- بحث الصور من أكثر من مصدر ---
def search_and_download_image(query: str, temp_dir: str) -> str:
    """
    دالة مبسطة للبحث عن صورة وتحميلها وإرجاع اسم الملف المحلي فقط
    """
    try:
        # البحث عن صورة من المصادر المختلفة
        img_url = None
        img_url = search_wikimedia(query)
        if not img_url:
            img_url = search_unsplash(query)
        if not img_url:
            img_url = search_pexels(query)
        
        if not img_url:
            return ""
        
        # تحميل الصورة
        response = requests.get(img_url, timeout=15)
        response.raise_for_status()
        
        # حفظ الصورة بمجلد مؤقت
        ext = img_url.split('.')[-1].lower()
        if ext not in ['jpg', 'jpeg', 'png', 'gif']:
            ext = 'jpg'
            
        filename = f"{hashlib.md5(query.encode()).hexdigest()}.{ext}"
        filepath = os.path.join(temp_dir, filename)
        
        with open(filepath, 'wb') as f:
            f.write(response.content)
            
        return filename  # إرجاع اسم الملف فقط (ليس المسار الكامل)
        
    except Exception as e:
        logging.error(f"❌ فشل في تحميل الصورة: {e}")
        return ""


def generate_special_anki_cards_from_text(content: str, major: str = "General", user_id: int = 0, num_cards: int = 15) -> tuple:
    for attempt in range(3):  # تجربة حتى 3 مرات
        prompt = f"""
You are an AI assistant specialized in creating study flashcards.

🎯 Task:
Extract the most important {num_cards} points from the following content, and convert each into an **Anki-style flashcard**.

🔹 Rules:
- Each flashcard must include:
  - "front": a short question or hint.
  - "back": the detailed answer or explanation.
  - "tag": (optional) topic label like Grammar, Biology, Logic, etc.
  - "image_hint": (optional) a short description of an image that would help illustrate the card (only if relevant).
- The front must be phrased to encourage recall (e.g. "What is...", "Define...", "How does...").
- Don't use Markdown, just clean plain text.
- Keep the cards diverse and helpful.
- Output must be a valid JSON **object** with two keys: "title" and "cards".

🚫 Important:
- Do NOT generate multiple choice or true/false questions.
- Only generate flashcards suitable for Anki with a front and a back.
- The flashcards must be written in the same language as the input content. If the content is in Arabic, answer in Arabic. If English, answer in English.

📘 Content to process (field: {major}):
{content}

✅ Example output format:
{{
  "title": "Basics of Organic Chemistry",
  "cards": [
    {{
      "front": "What is the function of mitochondria?",
      "back": "It is the powerhouse of the cell.",
      "tag": "Biology",
      "image_hint": "microscopic image of mitochondria"
    }},
    {{
      "front": "ما هي الاستعارة؟",
      "back": "الاستعارة هي استخدام الكلمة في غير معناها الحقيقي لعلاقة مع قرينة مانعة.",
      "tag": "Literature",
      "image_hint": ""
    }}
  ]
}}
"""

        # توليد الاستجابة
        if user_id == ADMIN_ID or can_generate(user_id):
            raw_output = generate_smart_response(prompt)
        else:
            raw_output = generate_gemini_response(prompt)

        # استخراج JSON
        clean_json = extract_json_from_string(raw_output)

        try:
            data = json.loads(clean_json)
            title = data.get("title", "بطاقات تعليمية")
            card_list = data.get("cards", [])

            cards = []
            for item in card_list:
                front = item.get("front") or item.get("question", "")
                back = item.get("back") or item.get("answer", "")
                tag = item.get("tag", "")
                image_hint = item.get("image_hint", "").strip()

                # التنظيف الأساسي
                front = str(front).strip()
                back = str(back).strip()
                tag = str(tag).strip()

                if front and back:  # بطاقة صالحة
                    cards.append({
                        "front": front,
                        "back": back,
                        "tag": tag,
                        "image_hint": image_hint  # سيتم معالجة الصور لاحقاً
                    })

            if len(cards) >= 5:
                return cards, title

        except json.JSONDecodeError as e:
            logging.error(f"❌ فشل في تحليل JSON للبطاقات: {e}")
            continue

    return [], "بطاقات تعليمية"

# -------------------------------------------------------------------
#                 games
# -------------------------------------------------------------------
import random

topics = [
    "حياة الطالب", "تخطيط السفر", "مشاريع جماعية", "مقابلات العمل",
    "الضغط الزمني", "مواقف عاطفية", "استخدام التكنولوجيا", "قرارات مالية",
    "صراعات الفريق", "تحديد الأهداف"
]

def generate_vocabulary_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "vocab")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""  
You are an AI vocabulary quiz creator.  
Generate one vocabulary question for a student majoring in {major}.
- Vocabulary should be relevant to real life or academic use and not an uncommon Vocabulary.
- Show the meaning of an English word in English 
- Provide 4 English words as options  
- Only ONE option should be correct.  
- Don't explain anything. Just give raw JSON.

Example:
{{
  "question": "Question",
  "options": ["Option", "Option", "Option", "Option"],
  "correct_index": 0
}}

Use this seed to diversify the question: {rand}
❌ Avoid repeating or paraphrasing these questions:
{recent_prompt}
"""
    q = generate_game(prompt)

    # حفظ السؤال الجديد
    add_recent_question(user_id, "speed", q["question"])
    return q

def generate_speed_challenge(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "speed")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""
You are a quiz bot.

Generate a **fun, fast-answer quiz** for a student in {major}.

Requirements:
- The question must be in English.
- The 4 options must be in English.
- Use fun and fast general knowledge topics (e.g. logic, daily life trivia, or language puzzles). Avoid repeating the same categories.
- Keep it simple and not too academic.
- Return raw JSON only.
- No explanation.
- Use this seed to increase randomness: {rand}
❌ Avoid repeating or paraphrasing these questions:
{recent_prompt}

Example output:
{{
  "question": "Question?",
  "options": ["Option", "Option", "Option", "Option"],
  "correct_index": 0
}}
"""
    q = generate_game(prompt, translate_question=True)

    # حفظ السؤال الجديد
    add_recent_question(user_id, "speed", q["question"])
    return q
    

# ★ لعبة الاخطاء الشائعة
def generate_common_mistakes_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "mistakes")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""
You are an educational game generator.

Your task:
- Generate a multiple-choice question highlighting a **common mistake** in the field of {major}.
- The question must be in English.
- The **options must be in English**.
- Provide **4 options** only, with one correct.
- Don't explain.
- Return only raw JSON.

❌ Avoid repeating or paraphrasing these questions:
{recent_prompt}
Use this random seed to diversify the question: {rand}

Example output:
{{
  "question": "Which sentence is grammatically incorrect?",
  "options": ["He go to school every day.", "She plays the piano.", "They are studying now.", "I have finished my homework."],
  "correct_index": 0
}}
"""
    q = generate_game(prompt, translate_question=True)

    # حفظ السؤال الجديد
    add_recent_question(user_id, "speed", q["question"])
    return q


def generate_inference_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "inference")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    random_topic = random.choice(topics)
    prompt = f"""
You are an AI-powered life skills test creator.

Generate a **new and unique** question that develops one of the following skills:  
- Critical thinking  
- Emotional intelligence  
- Time management  
- Self-awareness  
- Decision making  
- Problem solving  
- Logic  
- Pattern recognition  
- Mental map understanding  

🔹 **Requirements**:  
- Write the **question in Arabic**  
- Write **all options in Arabic**  
- Use a realistic scenario or student-life context related to: **{random_topic}**  
- Provide **exactly 4 options**, with **one correct answer**  
- **Never repeat** past examples or add explanations  
- Make the question **engaging and clever**  
- Incorporate variability using this random number: **{rand}**  
- the options should be as short as possible but understandable
❌ Avoid repeating or paraphrasing these questions:
{recent_prompt}
🔸 Return **JSON-only output** (no additional text).  

Example (Johnson’s format):  
{{
  "question": "Question",  
  "options": ["Options", "Option", "Option", "Option"],  
  "correct_index": 2  
}}  
"""
    q = generate_game(prompt, translate_question=True)

    # حفظ السؤال الجديد
    add_recent_question(user_id, "speed", q["question"])
    return q
    
# ----------------------------------
# ------------- inference review -------------------------------------------------------------------


def review_inference_question_with_ai(question_text: str, options: list[str], correct_index: int) -> bool:
    prompt = f"""
You are an AI educational assistant.

A student submitted the following inference question. Review it and decide if it's valid:
- Is the question clear and meaningful?
- Are the 4 options distinct and related to the question?
- Is there **one and only one correct answer**?

Respond only with YES or NO.

Question: {question_text}
Options: {options}
Correct index: {correct_index}
"""
    response = generate_smart_response(prompt).strip().lower()
    return "yes" in response

# عدّل هذه الدالة
def send_quiz_to_user(chat_id, quiz_data, message_id=None): # message_id يصبح اختيارياً
    markup = InlineKeyboardMarkup()
    quiz_url = f"{WEBHOOK_URL}/quiz/{quiz_data['db_id']}"
    btn = InlineKeyboardButton("🚀 فتح الاختبار", url=quiz_url)
    markup.add(btn)
    
    message = f"""
    🏆 تم إنشاء اختبارك الطبي بنجاح!
    العنوان: {quiz_data['title']}
    عدد الأسئلة: {len(quiz_data['questions'])}
    """
    # احذف الشرط وأرسل رسالة جديدة دائمًا
    bot.send_message(
        chat_id,
        message,
        reply_markup=markup
    )
    
def process_pending_inference_questions():
    cursor.execute("SELECT id, question, options, correct_index FROM inference_questions WHERE approved = 0")
    pending = cursor.fetchall()

    for row in pending:
        qid, qtext, options_json, correct_index = row
        try:
            options = json.loads(options_json)
        except:
            continue  # تجاهل الأسئلة ذات التنسيق الخاطئ

        if review_inference_question_with_ai(qtext, options, correct_index):
            cursor.execute("UPDATE inference_questions SET approved = 1 WHERE id = ?", (qid,))
        else:
            cursor.execute("DELETE FROM inference_questions WHERE id = ?", (qid,))

    conn.commit()


import sqlite3
import time
import json
import uuid
import threading
from datetime import datetime
from telebot import types
import uuid

from datetime import datetime

def log_quiz_share(quiz_code, shared_by_user_id, shared_by_name):
    conn = sqlite3.connect("quiz_users.db")
    c = conn.cursor()

    shared_at = datetime.now().isoformat()

    c.execute("""
        INSERT INTO quiz_shares (quiz_code, shared_by_user_id, shared_by_name, shared_at)
        VALUES (?, ?, ?, ?)
    """, (quiz_code, shared_by_user_id, shared_by_name, shared_at))

    conn.commit()
    conn.close()



# توليد كود فريد مع التحقق من التكرار
def generate_unique_quiz_code():
    while True:
        code = f"QC_{uuid.uuid4().hex[:6]}"
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()
        c.execute("SELECT 1 FROM user_quizzes WHERE quiz_code = ?", (code,))
        if not c.fetchone():
            conn.close()
            return code
        conn.close()

# تخزين الاختبار (نسخة محسنة)
def store_quiz(user_id, quizzes, bot):
    try:
    
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()

        quiz_code = generate_unique_quiz_code()

        # جلب اسم المالك
        try:
            owner_chat = bot.get_chat(user_id)
            owner_name = owner_chat.first_name or owner_chat.username or f"user_{user_id}"
        except Exception:
            owner_name = "صديقك"

        created_at = datetime.now().isoformat()

        c.execute("""
            INSERT INTO user_quizzes (user_id, quiz_data, quiz_code, created_at, is_active, owner_name)
            VALUES (?, ?, ?, ?, 1, ?)
        """, (user_id, json.dumps(quizzes), quiz_code, created_at, owner_name))

        conn.commit()
        conn.close()

        return quiz_code
    except Exception as e:
        print(f"Error storing quiz: {e}")
        return None



def start_quiz(chat_id, quiz_code, bot):
    if quiz_code == "sample":
        # استعلام من جدول العينة
        cursor.execute("SELECT quiz_data FROM sample_quizzes WHERE quiz_code = ?", (quiz_code,))
        row = cursor.fetchone()
        if not row:
            return False
        quiz_data_json = row[0]
    else:
        # استعلام من جدول المستخدمين
        cursor.execute("SELECT quiz_data FROM user_quizzes WHERE quiz_code = ?", (quiz_code,))
        row = cursor.fetchone()
        if not row:
            return False
        quiz_data_json = row[0]

    # تحويل JSON إلى بيانات، ثم متابعة تشغيل الاختبار
    quiz_data = json.loads(quiz_data_json)
    # تابع عرض الأسئلة والتعامل مع المستخدم هنا...
    # ...

    return True
    
    
    
logger = logging.getLogger(__name__)


# نظام إدارة الحالة المحسن
class QuizManager:

    def __init__(self):
        self.active_quizzes = {}

    def start_quiz(self, chat_id, quiz_code, bot2, message_id=None):
        try:
            # ----- جلب الأسئلة حسب نوع الكود -----
            if quiz_code == "sample":
                # افتح اتصالًا بقاعدة quiz_users.db حيث مخزن sample_quizzes
                conn = sqlite3.connect('quiz_users.db')
                cur = conn.cursor()
                cur.execute("SELECT quiz_data FROM sample_quizzes WHERE quiz_code = ?", (quiz_code,))
                row = cur.fetchone()
                conn.close()

                if not row:
                    logger.info("sample quiz not found in sample_quizzes")
                    return False

                # row[0] يفترض أنه JSON نصي
                quizzes_raw = json.loads(row[0]) if isinstance(row[0], str) else row[0]
                # sample عادة لا يملك owner
                user_id = None
                score = 0
                total = len(quizzes_raw)

                shared_by_name = "TestGenie"   # اسم افتراضي للعرض
                owner_name = "المرسل"

            else:
                # استعلام من user_quizzes
                conn = sqlite3.connect("quiz_users.db")
                c = conn.cursor()
                c.execute("""
                    SELECT user_id, score, total, quiz_data, owner_name
                    FROM user_quizzes
                    WHERE quiz_code = ?
                """, (quiz_code,))
                row = c.fetchone()
                conn.close()

                if not row:
                    logger.info("user quiz not found for code %s", quiz_code)
                    return False

                # تفكيك row بأمان
                user_id = row[0]
                score = row[1] or 0
                total = row[2] or 0
                quiz_data_field = row[3]
                owner_name = row[4] or "المالك"

                # quiz_data_field يمكن أن يكون نص JSON أو لقائمة بالفعل
                quizzes_raw = json.loads(quiz_data_field) if isinstance(quiz_data_field, str) else quiz_data_field
                shared_by_name = owner_name

            # ----- تحويل أي بنية (dict أو list) إلى قائمة موحدة من القواميس -----
            formatted_quizzes = []
            for q in quizzes_raw:
                if isinstance(q, dict):
                # بنية القاموس المتوقعة
                    question = q.get('question')
                    options = q.get('options')
                    correct_idx = q.get('correct_index')
                    explanation = q.get('explanation', '')
                elif isinstance(q, (list, tuple)) and len(q) >= 4:
                    # بنية قائمة قد تكون [question, options, correct_idx, explanation]
                    question, options, correct_idx, explanation = q[0], q[1], q[2], q[3]
                else:
                    logger.warning("Unknown question format: %r", q)
                    continue

                # تحقق من صحة الخيارات
                if not (isinstance(options, list) and len(options) == 4):
                    logger.warning("Skipping question with invalid options: %r", question)
                    continue
                if not isinstance(correct_idx, int) or not (0 <= correct_idx < 4):
                    logger.warning("Skipping question with invalid correct_index: %r", question)
                    continue

                formatted_quizzes.append({
                    'question': question,
                    'options': options,
                    'answer': options[correct_idx],
                    'explanation': explanation
                })

            if not formatted_quizzes:
                logger.info("No valid questions found for quiz %s", quiz_code)
                return False

            # ----- احفظ الحالة النشطة للامتحان -----
            self.active_quizzes[chat_id] = {
                'quizzes': formatted_quizzes,
                'current_index': 0,
                'score': 0,
                'quiz_code': quiz_code,
                'start_time': datetime.now(),
                'owner_id': user_id,
                'owner_score': score,
                'owner_total': total
            }

            # ----- رسالة تقديمية -----
            estimated_time = round((total or len(formatted_quizzes)) * 0.5)  # نصف دقيقة لكل سؤال تقريبًا

            if user_id is not None and chat_id == user_id:
                msg = f"تم توليد {len(formatted_quizzes)} سؤالًا، استعد للاختبار."
            else:
                percent = round((score / total) * 100) if total else 0
                msg = (
                    f"🎯 <b>{shared_by_name} أرسل لك تحديًا!</b> 🤝\n\n"
                    f"📋 <b>عدد الأسئلة:</b> {total or len(formatted_quizzes)}\n"
                    f"🏆 <b>نتيجة {owner_name}:</b> {score}/{total if total else len(formatted_quizzes)} — ({percent}%)\n"
                    f"⏳ <b>الوقت المقدر:</b> حوالي {estimated_time} دقيقة\n\n"
                    f"🔥 <b>هل تستطيع التفوق عليه؟</b>"
                )

            # ارسل أو عدّل رسالة التحميل
            if message_id:
                try:
                    bot2.edit_message_text(chat_id=chat_id, message_id=message_id, text=msg, parse_mode="HTML")
                except Exception as e:
                    logger.exception("Failed to edit message: %s", e)
            else:
                bot2.send_message(chat_id, msg, parse_mode="HTML")
                time.sleep(1)

            # أرسل السؤال الأول
            self.send_question(chat_id, bot2)
            return True

        except Exception as e:
            logger.exception("start_quiz failed: %s", e)
            return False
    
    
    def get_quiz_info(self, quiz_code):
        """الحصول على معلومات الاختبار من قاعدة البيانات"""
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()
        c.execute("SELECT user_id, quiz_data FROM user_quizzes WHERE quiz_code = ?", (quiz_code,))
        result = c.fetchone()
        conn.close()
        
        if not result:
            return None
            
        return {
            'user_id': result[0],
            'quizzes': json.loads(result[1])
        }
    
    def can_access_quiz(self, user_id, quiz_code):
        """التحقق من صلاحية الوصول للاختبار"""
        info = self.get_quiz_info(quiz_code)
        if not info:
            return False
            
            
    def send_question(self, chat_id, bot2):
        state = self.active_quizzes.get(chat_id)
        if not state:
            return
            
        quiz = state['quizzes'][state['current_index']]
        
        try:
            poll = bot2.send_poll(
                chat_id=chat_id,
                question=quiz['question'],
                options=quiz['options'],
                type='quiz',
                correct_option_id=quiz['options'].index(quiz['answer']),
                explanation=quiz['explanation'],
                is_anonymous=False,
                open_period=30
            )
            
            state['last_poll_id'] = poll.message_id
        except Exception as e:
            print(f"Error sending poll: {e}")
            self.handle_quiz_end(chat_id, bot2, error=True)


    def handle_answer(self, poll_answer, bot2):
        chat_id = poll_answer.user.id
        state = self.active_quizzes.get(chat_id)
        if not state:
            return
            
        current_quiz = state['quizzes'][state['current_index']]
        is_correct = poll_answer.option_ids[0] == current_quiz['options'].index(current_quiz['answer'])
        
        if is_correct:
            state['score'] += 1
            feedback = "✅ إجابة صحيحة!"
        else:
            feedback = f"❌ إجابة خاطئة! الإجابة الصحيحة هي: {current_quiz['answer']}"
            
        if current_quiz['explanation']:
            feedback += f"\n\n💡 التفسير: {current_quiz['explanation']}"
            
        bot2.send_message(chat_id, feedback)
        
        # الانتقال للسؤال التالي
        state['current_index'] += 1
        if state['current_index'] < len(state['quizzes']):
            self.send_question(chat_id, bot2)
        else:
            self.handle_quiz_end(chat_id, bot2)

    def handle_quiz_end(self, chat_id, bot2, error=False):
        state = self.active_quizzes.pop(chat_id, None)
        if not state:
            return
            
        if error:
            bot2.send_message(chat_id, "⚠️ حدث خطأ في الاختبار. يرجى المحاولة لاحقًا")
            return
            
        total = len(state['quizzes'])
        score = state['score']
        quiz_code = state['quiz_code']

                # حفظ النتيجة في DB
        conn = sqlite3.connect("quiz_users.db")
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE user_quizzes
            SET score = ?, total = ?, timestamp = ?
            WHERE quiz_code = ?
        """, (score, total, datetime.now().isoformat(), state['quiz_code']))
        conn.commit()
        conn.close()

        keyboard = types.InlineKeyboardMarkup()
        keyboard.add(types.InlineKeyboardButton("🔄 إعادة الاختبار", callback_data=f"retry:{quiz_code}"))
        keyboard.add(types.InlineKeyboardButton("📤 مشاركة الاختبار", callback_data=f"share_quiz:{quiz_code}"))
        keyboard.add(types.InlineKeyboardButton("➡️ العودة الى TestGenie ✨", url="https://t.me/Oiuhelper_bot"))
        
        end_msg = bot2.send_message(
            chat_id,
            f"🎉 انتهى الاختبار!\n\nنتيجتك: {score}/{total}\n\nماذا تريد أن تفعل الآن؟",
            reply_markup=keyboard
        )

# تهيئة المدير
quiz_manager = QuizManager()


@bot2.poll_answer_handler()
def handle_poll_answer(poll_answer):
    quiz_manager.handle_answer(poll_answer, bot2)


# دالة إرسال الاختبارات (النسخة المحسنة)
def send_quizzes(chat_id, quizzes, message_id=None):
    try:
        # إرسال رسالة التحميل
        msg = bot2.send_message(chat_id, "⏳ جاري تحضير الاختبار...")
        
        # تخزين الاختبار
        quiz_code = store_quiz(chat_id, quizzes)
        if not quiz_code:
            raise Exception("Failed to store quiz")
        
        # بدء الاختبار
        if quiz_manager.start_quiz(chat_id, quiz_code, bot2):
            bot.delete_message(chat_id, msg.message_id)
        else:
            bot2.edit_message_text("❌ فشل بدء الاختبار", chat_id, msg.message_id)
            
    except Exception as e:
        print(f"Error in send_quizzes: {e}")
        bot2.send_message(chat_id, "حدث خطأ أثناء بدء الاختبار")







import logging
import threading

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

state_lock = threading.Lock()
user_states = {}  # global



# -------------------------------------------------------------------
#                  Telegram Bot Handlers
# -------------------------------------------------------------------


@bot2.message_handler(commands=['start'])
def unified_start_handler(message):
    if message.chat.type != "private":
        return
    
    chat_id = message.chat.id
    args = message.text.split()
    
    if len(args) == 1:
        bot2.send_message(chat_id, """
        👋 أهلاً بك في Quizzy! 😊
        أنا هنا لمساعدتك في إجراء إختبارات سريعة والاستمتاع بالتعلم.
        """)
        return
    
    # استخراج الكود
    param = args[1]
    quiz_code = param[5:] if param.startswith("quiz_") else param

    # حالة خاصة لاختبار sample
    if quiz_code == "sample":
        loading_msg = bot2.send_message(chat_id, "🧠 جاري تحميل الاختبار التجريبي...")
        if not quiz_manager.start_quiz(chat_id, quiz_code, bot2, loading_msg.message_id):
            bot2.edit_message_text(
                chat_id=chat_id,
                message_id=loading_msg.message_id,
                text="❌ لم يتم العثور على الاختبار التجريبي."
            )
        return
    
    # أي كود اختبار آخر
    loading_msg = bot2.send_message(chat_id, "🧠 جاري تحميل الاختبار...")
    if not quiz_manager.start_quiz(chat_id, quiz_code, bot2, loading_msg.message_id):
        bot2.edit_message_text(
            chat_id=chat_id,
            message_id=loading_msg.message_id,
            text="❌ لم يتم العثور على هذا الاختبار أو قد تكون صلاحيته انتهت."
            )
    # إذا كان كود مختلف، تتصرف حسب الكود عادي
    # ... باقي المعالجة ...
    
@bot2.callback_query_handler(func=lambda c: True)
def handle_main_menu(c):
    
    if c.message.chat.type != "private":
        return
    uid = c.from_user.id
    data = c.data
    chat_id = c.message.chat.id
    message_id = c.message.message_id

    if data.startswith("retry:"):
        quiz_code = data[6:]
        quiz_manager.start_quiz(chat_id, quiz_code, bot2)
            
        
    elif data.startswith("share_quiz:"):
        quiz_code = data[6:]
        chat_id = c.message.chat.id  # ← تأكد من تعيين chat_id هنا

        try:
            user_chat = bot2.get_chat(uid)
            shared_by_name = user_chat.first_name or user_chat.username or f"user_{uid}"
        except Exception:
            shared_by_name = "صديقك"

        log_quiz_share(quiz_code, uid, shared_by_name)
        file_path = user_files[uid]
        
        share_link = f"https://t.me/QuizzyAI_bot?start=quiz_{quiz_code}"
        
        msg_text_share = f"""📢 {shared_by_name} أرسل لك هذا الاختبار!  

📂 الملف: {msg.document.file_name}

جربه واختبر معلوماتك 👇  
{share_link}
"""
        msg_text = f"""<b>🎉 شارك هذا الاختبار مع زملائك!</b>

    انسخ الرابط أدناه أو اضغط لفتحه مباشرة:
    🔗 <a href="{share_link}">{share_link}</a>

    📝 عند فتح الرابط، سيبدأ الاختبار تلقائيًا بإذن الله.  
    📢 بمشاركتك هذا الاختبار قد يصير عامًا.
    """

        keyboard = types.InlineKeyboardMarkup()
        keyboard.add(
            types.InlineKeyboardButton("🔗 نسخ الرابط", switch_inline_query=msg_text_share),
            types.InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home")
        )

        bot2.edit_message_text(msg_text, chat_id=chat_id, message_id=message_id, parse_mode="HTML", reply_markup=keyboard)


@bot.message_handler(commands=['start'])
def unified_start_handler(message):
    # ✅ تجاهل الرسائل في المجموعات
    if message.chat.type != "private":
        return

    chat_id = message.chat.id
    args = message.text.split()
    uid = message.from_user.id
    
    if not can_generate(uid):
        add_external_user(uid)
    
    # ✅ إذا وُجد باراميتر (مثل quiz_ab12cd أو anki_sample)
    if len(args) > 1:
        param = args[1]

        # ✅ إذا كان باراميتر anki_sample
        if param == "anki_sample":
            user_states[chat_id] = "awaiting_anki_file_ai"  # حفظ الحالة
            bot.send_message(
                "📝 دعنا نبدأ بإنشاء **ملف بطاقاتك الأول**!\n"
                "📂 أرسل ملف **PDF** أو **DOCX** أو **PPTX**، أو حتى نصًا مباشرًا 📜.\n"
                "سيتم توليد ملف **أنكي** مخصص لك تلقائيًا 🎯",
                chat_id=chat_id,
                parse_mode="Markdown"
            )

            return

        # ✅ معالجة روابط المشاركة مثل: ?start=quiz_ab12cd
        quiz_code = param[5:] if param.startswith("quiz_") else param

        loading_msg = bot.send_message(chat_id, "🧠 جاري تحميل الاختبار...")

        # ✅ محاولة بدء الاختبار
        if not quiz_manager.start_quiz(chat_id, quiz_code, bot):
            bot.edit_message_text(
                chat_id=chat_id,
                message_id=loading_msg.message_id,
                text="❌ لم يتم العثور على هذا الاختبار أو انتهت صلاحيته."
            )
        return

    # ✅ إذا لم يوجد باراميتر → عرض القائمة الرئيسية
    send_main_menu(chat_id)


def send_main_menu(chat_id, message_id=None):
    keyboard = InlineKeyboardMarkup(row_width=2)
    buttons = [
        InlineKeyboardButton("📝 توليد اختبار", callback_data="go_generate"),
        InlineKeyboardButton("📚 مراجعة سريعة", callback_data="soon_review"),
        InlineKeyboardButton("📄 ملخص PDF", callback_data="soon_summary"),
        InlineKeyboardButton("🧠 بطاقات Anki", callback_data="anki"),
        InlineKeyboardButton("🎮 ألعاب تعليمية", callback_data="go_games"),
        InlineKeyboardButton("⚙️ حسابي", callback_data="go_account_settings"),
    ]
    keyboard.add(*buttons)
    keyboard.add(InlineKeyboardButton("➕ أضفني إلى مجموعة", url=f"https://t.me/{bot.get_me().username}?startgroup=true"))

    text = (
        "👋 <b>أهلاً بك في TestGenie!</b> ✨\n\n"
        "🎯 أدوات تعليمية ذكية بين يديك:\n"
        "- اختبارات من ملفاتك\n"
        "- بطاقات مراجعة (Anki)\n"
        "- ملخصات PDF/Word <i>(قريباً)</i>\n"
        "- ألعاب تعليمية ممتعة\n\n"
        "📌 كل ما تحتاجه لتتعلّم بذكاء... بين يديك الآن.\n\n"
        "👇 اختر ما يناسبك وابدأ الآن:"
    )

    if message_id:
        bot.edit_message_text(
            text,
            chat_id=chat_id,
            message_id=message_id,
            reply_markup=keyboard,
            parse_mode="HTML"
        )
    else:
        bot.send_message(
            chat_id,
            text,
            reply_markup=keyboard,
            parse_mode="HTML"
        )



@bot.callback_query_handler(func=lambda call: call.data.startswith("rate_"))
def handle_rating(call):
    uid = call.from_user.id
    username = call.from_user.username or "مستخدم"
    rating = call.data.replace("rate_", "")

    if rating == "ignore":
        bot.answer_callback_query(call.id, "✅ تم تجاهل التقييم.")
        return

    # إشعار الأدمن عبر bot3
    bot3.send_message(
        ADMIN_ID,
        f"⭐ تقييم جديد من @{username} (UID: {uid})\n\nالتقييم: {rating} نجوم"
    )

    # رسالة للمستخدم
    bot.send_message(
        call.message.chat.id,
        f"💌 شكراً لك @{username}! تم إرسال ملاحظاتك ({rating}⭐) إلى فريق TestGenie 🙏"
    )

    bot.answer_callback_query(call.id, "✅ تم تسجيل تقييمك، شكراً لك!")


@bot.callback_query_handler(func=lambda c: True)
def handle_main_menu(c):
    try:
        bot.answer_callback_query(c.id)
    except:
        pass

    if c.message.chat.type != "private":
        return
    try: 
    
        data = c.data
        chat_id = c.message.chat.id
        message_id = c.message.message_id
        uid = c.from_user.id
        logging.info("Callback received: uid=%s data=%s", uid, data)


        # ردود خاطئة عشوائية تظهر للمستخدم
        wrong_responses = [
            "❌ خطأ! جرب مجددًا 😉\n✅ الصحيح: {correct}",
            "🚫 للأسف، ليست الصحيحة!\n✅ الجواب: {correct}",
            "😅 ليست الإجابة الصحيحة، الجواب هو: {correct}",
            "❌ لا، حاول مرة أخرى!\n✔️ الصحيح هو: {correct}"
        ]




    # ---------- صفحة الفئات الأولى ----------
        if data == "go_generate":
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("🩺 الطب والصحة", "category:health:page1"),
                ("🛠️ الهندسة", "category:engineering:page1"),
                ("💻 الحاسوب", "category:computer:page1"),
                ("📊 الإدارة", "category:business:page1"),
                ("🗣️ اللغات", "category:languages:page1"),
            ]
            # بناء الصفوف بأمان (لا نمرر None)
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            keyboard.add(InlineKeyboardButton("❓ تخصص آخر", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("➡️ المزيد من التخصصات", callback_data="go_next"))
            keyboard.add(InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "🎯 *اختر مجال تخصصك* (1/2)\n\nحدد الفئة الأقرب لتخصصك من القائمة:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
            )
            except Exception as e:
                print("خطأ في عرض الصفحة الأولى:", e)
                # محاولة إرسال رسالة جديدة
                try:
                    bot.send_message(chat_id, "🎯 اختر مجال تخصصك (1/2).", reply_markup=keyboard)
                except Exception as e2:
                    print("فشل الإرسال:", e2)
            return

    # ---------- صفحة الفئات الثانية ----------
        if data == "go_next":
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("📿 العلوم الإسلامية", "category:islamic:page2"),
                ("⚖️ القانون والسياسة", "category:law:page2"),
                ("🧪 العلوم الطبيعية", "category:science:page2"),
                ("🎨 الفنون", "category:arts:page2"),
                ("👩‍🏫 التربية", "category:education:page2"),
            ]
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            keyboard.add(InlineKeyboardButton("❓ تخصص آخر", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("⬅️ العودة", callback_data="go_prev"))
            keyboard.add(InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "🎓 *المزيد من التخصصات* (2/2)\n\nاختر من الفئات التالية:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("خطأ في عرض الصفحة الثانية:", e)
            return

    # ---------- العودة للصفحة الأولى ----------
        if data == "go_prev":
        # إعادة توجيه إلى نفس تنفيذ go_generate
        # نستدعي نفس الكود أو نعيد التوجيه:
        # لإعادة الاستخدام يمكن عمل دالة build_go_generate_keyboard()
            bot.answer_callback_query(c.id)
        # هنا نعيد نفس الكود أعلاه أو ننادي تنفيذ go_generate:
        # أسهل: نعيد تعيين data ونكولّد الصفحة:
        # ولكن لتبسيط، نعيد واجهة go_generate مباشرة:
        # (يمكنك استدعاء رسالة موجودة مسبقًا)
        # نعيد نفس لوحة go_generate كما في الأعلى:
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("🩺 الطب والصحة", "category:health:page1"),
                ("🛠️ الهندسة", "category:engineering:page1"),
                ("💻 الحاسوب", "category:computer:page1"),
                ("📊 الإدارة", "category:business:page1"),
                ("🗣️ اللغات", "category:languages:page1"),
            ]
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)
            keyboard.add(InlineKeyboardButton("❓ تخصص آخر", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("➡️ المزيد من التخصصات", callback_data="go_next"))
            keyboard.add(InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "🎯 *اختر مجال تخصصك* (1/2)\n\nحدد الفئة الأقرب لتخصصك من القائمة:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("خطأ في عرض الصفحة الأولى (go_prev):", e)
            return

    # ---------- معالجة اختيار الفئة (آمنة) ----------
        if data.startswith("category:"):
            # الصيغة المتوقعة: "category:<key>:<page>"
            parts = data.split(":")
            if len(parts) != 3:
                print("callback category غير متوافق:", data)
                return
            _, cat_key, page = parts  # cat_key مثلاً 'health' ، page مثل 'page1'

            # خريطة التخصصات الفرعية (استخدم المفاتيح الإنجليزية)
            SUBS = {
                "health": [
                    ("🧬 الطب البشري", "major:طب_بشري"),
                    ("💊 الصيدلة", "major:صيدلة"),
                    ("🏥 التمريض", "major:تمريض"),
                    ("🔬 علوم المختبرات", "major:علوم_مختبرات"),
                    ("🦷 طب الأسنان", "major:طب_أسنان"),
                ],
                "engineering": [
                    ("⚙️ الميكانيكا", "major:هندسة_ميكانيكية"),
                    ("🧪 الكيميائية", "major:هندسة_كيميائية"),
                    ("💡 الكهربائية", "major:هندسة_كهربائية"),
                    ("🏗️ المدنية", "major:هندسة_مدنية"),
                    ("🔧 البرمجيات", "major:هندسة_برمجيات"),
                    ("📡 الاتصالات", "major:هندسة_اتصالات"),
                ],
                "computer": [
                    ("💻 علوم الحاسوب", "major:علوم_حاسوب"),
                    ("📱 تطوير البرمجيات", "major:تطوير_برمجيات"),
                    ("🔒 أمن المعلومات", "major:امن_معلومات"),
                    ("🤖 الذكاء الاصطناعي", "major:ذكاء_اصطناعي"),
                    ("📊 علم البيانات", "major:علم_البيانات"),
                ],
                "business": [
                    ("📈 إدارة الأعمال", "major:ادارة_اعمال"),
                    ("💹 الاقتصاد", "major:اقتصاد"),
                    ("📊 المحاسبة", "major:محاسبة"),
                    ("🧮 التسويق", "major:تسويق"),
                    ("🏦 التمويل", "major:تمويل"),
                ],
                "languages": [
                    ("🌎 الإنجليزية", "major:لغة_انجليزية"),
                    ("🇫🇷 الفرنسية", "major:لغة_فرنسية"),
                    ("🇸🇦 العربية", "major:لغة_عربية"),
                    ("📚 الترجمة", "major:ترجمة"),
                    ("🇩🇪 الألمانية", "major:لغة_ألمانية"),
                ],
                "islamic": [
                    ("📜 الفقه", "major:فقه"),
                    ("💡 العقيدة", "major:عقيدة"),
                    ("📖 التفسير", "major:تفسير"),
                    ("🕌 الدراسات الإسلامية", "major:دراسات_اسلامية"),
                    ("🌙 السيرة النبوية", "major:سيرة_نبوية"),
                ],
                "law": [
                    ("📜 القانون الدولي", "major:قانون_دولي"),
                    ("🏛️ القانون الوطني", "major:قانون_وطني"),
                    ("🗳️ العلوم السياسية", "major:علوم_سياسية"),
                    ("👮‍♂️ القانون الجنائي", "major:قانون_جنائي"),
                    ("⚖️ القانون التجاري", "major:قانون_تجاري"),
                ],
                "science": [
                    ("🧪 الفيزياء", "major:فيزياء"),
                    ("🔬 الكيمياء", "major:كيمياء"),
                    ("🔢 الرياضيات", "major:رياضيات"),
                    ("🌿 الأحياء", "major:احياء"),
                ],
                "arts": [
                    ("🎭 الفنون الأدائية", "major:فنون_ادائية"),
                    ("🖼️ الفنون البصرية", "major:فنون_بصرية"),
                    ("📚 الأدب", "major:ادب"),
                    ("🌏 التاريخ", "major:تاريخ"),
                    ("🎵 الموسيقى", "major:موسيقى"),
                ],
                "education": [
                    ("👩‍🏫 علم النفس التربوي", "major:علم_النفس_تربوي"),
                    ("📘 المناهج", "major:مناهج_تدريس"),
                    ("🧩 التربية الخاصة", "major:تربية_خاصة"),
                    ("📊 إدارة التعليم", "major:ادارة_تعليم"),
                    ("👶 الطفولة المبكرة", "major:طفولة_مبكرة"),
                ],
            }

            sub_btns = SUBS.get(cat_key, [])
            keyboard = InlineKeyboardMarkup(row_width=1)
            for label, cb in sub_btns:
                keyboard.add(InlineKeyboardButton(label, callback_data=cb))

            keyboard.add(InlineKeyboardButton("❓ تخصص آخر", callback_data="major_custom"))

            # زر الرجوع: إذا كانت الصفحة page1 نرجع الى go_generate وإلا نرجع الى go_next
            back_to = "go_generate" if page == "page1" else "go_next"
            keyboard.add(InlineKeyboardButton("⬅️ رجوع", callback_data=back_to))

            try:
                pretty_name = cat_key.replace("_", " ")
                bot.edit_message_text(
                    f"🎓 *{pretty_name}*\n\n👇 اختر تخصصك الدقيق من القائمة:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("خطأ في تعديل الرسالة (category):", e)
                try:
                    bot.send_message(chat_id, f"🎓 {pretty_name}\nاختر تخصصك:", reply_markup=keyboard)
                except Exception as e2:
                    print("فشل إرسال رسالة الفرعية:", e2)
            
            return



    # ---------- التخصص المخصص ----------
        if data == "major_custom":
            user_states[uid] = "awaiting_major"
            bot.edit_message_text("📝 أرسل اسم تخصصك (مثال: هندسة طيران، علم البيانات):", chat_id=chat_id, message_id=message_id)
            
            return



    # ----------------- المعالجة عند اختيار التخصص -----------------


        if data.startswith("major:"):
            major_key = data.split(":", 1)[1]  # 'طب_بشري' أو 'هندسة_ميكانيكية'

            # احفظ التخصص في DB (تأكد أن conn,cursor موجودان)
            try:
                cursor.execute("INSERT OR REPLACE INTO users(user_id, major) VALUES(?, ?)", (uid, major_key))
                conn.commit()
            except Exception as e:
                print("خطأ في حفظ التخصص:", e)

            # سلوك خاص لـ طب_بشري
            if major_key == "طب_بشري":
                kb2 = InlineKeyboardMarkup(row_width=2)
                kb2.add(InlineKeyboardButton("🧠 الوضع المتقدم", callback_data="advanced"), InlineKeyboardButton("📄 الوضع المباشر", callback_data="simple"))
                bot.edit_message_text(
                    "🧬 *نظام الاختبارات المتقدم - لطلبة الطب البشري*\n\n"
                    "يسرّنا إعلامك بأن نظامنا الجديد لتوليد الاختبارات يعتمد على الذكاء الاصطناعي المتقدم، "
                    "ويقوم بإعداد اختبارات *ذكية ومتوازنة* تحاكي هيكلة الامتحانات الفعلية من حيث التدرج والتنوع.\n\n"
                    "📌 اختر أحد الأوضاع التالية لبدء توليد اختبارك:",
                    chat_id=chat_id,
                    message_id=message_id,
                    reply_markup=kb2,
                    parse_mode="Markdown"
                )
            else:
                user_states[uid] = "awaiting_simple_test_file"

                sent_msg = bot.edit_message_text(
                    f"✅ تم تحديد تخصصك: *{major_key.replace('_', ' ')}*\nالآن أرسل ملف (PDF/DOCX/TXT) أو نصًا مباشرًا لتوليد اختبارك.",
                    chat_id=chat_id, message_id=message_id, parse_mode="Markdown"
            )
                return
            return

    # ... بقية المعالجات الأخرى ...


        if data == "anki":
            bot.answer_callback_query(c.id)
            choice_markup = types.InlineKeyboardMarkup()
            choice_markup.row(
                types.InlineKeyboardButton("📝 توليد يدوي", callback_data="manual_anki"),
                types.InlineKeyboardButton("🤖 توليد آلي", callback_data="ai_anki")
            )
            bot.edit_message_text(
                chat_id=chat_id,
                message_id=message_id,
                text="🔧 حدد طريقة إنشاء بطاقات Anki:",
                reply_markup=choice_markup
            )

        elif data == "manual_anki":
            with state_lock:
                user_states[uid] = "awaiting_anki_file_manual"
            bot.answer_callback_query(c.id, "✅ تم تفعيل الوضع اليدوي — أرسل البطائق كنص.")
            bot.edit_message_text(
                chat_id=chat_id, message_id=message_id,
                text=(
                    "✏️ *صياغة بطاقات Anki يدويًا* \n\n"
                    "يرجى إرسال البطاقات بنسق محدد كما يلي:\n\n"
                    "السؤال أو المصطلح.\n"
                    "الجواب أو التعريف.\n"
                    "تاق *(إختياري)*\n\n"
                    "مثال:\n"
                    "ما هو الذكاء الاصطناعي؟\n"
                    "هو قدرة الآلة على محاكاة التفكير البشري.\n"                
                    "#تقنية\n\n"
                    "من اكتشف الجاذبية؟\n"
                    "نيوتن\n\n"
                    "هذا النسق يساعد في توليد ملف Anki بسهولة ودقة. 💡"
                ),
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return


        elif data == "ai_anki":
            with state_lock:
                user_states[uid] = "awaiting_anki_file_ai"
            bot.answer_callback_query(c.id, "✅ تم تفعيل وضع Anki آلي — أرسل الملف الآن.")
            bot.edit_message_text(
                chat_id=chat_id, message_id=message_id,
                text=(
                    "✨ *توليد بطاقات Anki بواسطة الذكاء الإصطناعي*\n\n"
                    "ارفع ملفًا بصيغة PDF، Word، أو نص عادي، وسيقوم الذكاء الإصطناعي بتحليل المحتوى وتوليد بطاقات Anki متوافقة بشكل تلقائي.\n\n"
                    "ميزة فريدة:\n"
                    "- تحويل المحتوى إلى ملف Anki جاهز للاستخدام.\n"
                    "- يدعم مجموعة واسعة من المواضيع والمصادر.\n"
                    "- وفر الوقت والجهد باستخدام تقنية الذكاء الإصطناعي المتقدمة.\n\n"
                    "ابدأ الآن وأرسل ملفك لتحويله إلى ملف Anki مخصص! 📚"
                ),
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return


        elif data == "go_account_settings":
            bot.answer_callback_query(c.id)
            settings_keyboard = InlineKeyboardMarkup()
            settings_keyboard.add(
                InlineKeyboardButton("🎓 تغيير التخصص", callback_data="change_specialty"),
            )
            settings_keyboard.add(
                InlineKeyboardButton("📉 مستوى الإختبارات", callback_data="tests_level"),
            )
            settings_keyboard.add(
                InlineKeyboardButton("⬅️ رجوع", callback_data="go_back_home")
            )

            bot.edit_message_text(
                text="⚙️ *إعدادات الحساب*\n\n"
                    "يمكنك تخصيص تجربتك التعليمية هنا.\n"
                    "اختر ما ترغب بتعديله 👇",
                chat_id=c.message.chat.id,
                message_id=c.message.message_id,
                reply_markup=settings_keyboard,
                parse_mode="Markdown"
            )
        
        elif data == "go_games":
            raw = fetch_user_major(uid)

            if not row:
                user_states[uid] = "awaiting_major_for_games"
                bot.send_message(uid, "🧠 قبل أن نبدأ اللعب، أخبرنا بتخصصك:")
                return

            keyboard = InlineKeyboardMarkup(row_width=1)
            keyboard.add(
                InlineKeyboardButton("🔒 العب في الخاص", callback_data="game_private"),
                InlineKeyboardButton("👥 العب في المجموعة", switch_inline_query="game"),
                InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home")
            )
            bot.edit_message_text(
                "🎮 اختر طريقة اللعب:\n\n"
                "- 🔒 في الخاص (ألعاب شخصية حسب تخصصك)\n"
                "- 👥 في المجموعة (شارك الأصدقاء بالتحدي!)",
                chat_id=chat_id,
                message_id=message_id,
                reply_markup=keyboard
            )
    
        elif data == "go_back_home":
            # إعادة عرض واجهة البداية
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                InlineKeyboardButton("📝 توليد اختبار", callback_data="go_generate"),
                InlineKeyboardButton("📚 مراجعة سريعة", callback_data="soon_review"),
                InlineKeyboardButton("📄 ملخص PDF", callback_data="soon_summary"),
                InlineKeyboardButton("🧠 بطاقات Anki", callback_data="anki"),
                InlineKeyboardButton("🎮 ألعاب تعليمية", callback_data="go_games"),
                InlineKeyboardButton("⚙️ حسابي", callback_data="go_account_settings"),
            ]
            keyboard.add(*buttons)

            keyboard.add(InlineKeyboardButton("➕ أضفني إلى مجموعة", url=f"https://t.me/{bot.get_me().username}?startgroup=true"))

            bot.edit_message_text(
                "👋 أهلا بك في *TestGenie* ✨\n\n"
                "🎯 أدوات تعليمية ذكية بين يديك:\n"
                "- اختبارات من ملفاتك\n"
                "- بطاقات مراجعة (Anki)\n"
                "- ملخصات PDF/Word _(قريباً)_\n"
                "- ألعاب تعليمية\n\n"
                "📌 كل ما تحتاجه لتتعلّم بذكاء... بين يديك الآن.\n\n"
                "اختر ما يناسبك وابدأ الآن 👇",
                chat_id=chat_id,
                message_id=message_id,
                reply_markup=keyboard,
                parse_mode="Markdown"
            )
                
        

    # ----------------- تغيير التخصص (الواجهة الرئيسية للفئات) -----------------
        if data == "change_specialty":
            keyboard = InlineKeyboardMarkup(row_width=2)
            main_categories = [
                ("🩺 العلوم الصحية", "spec_category:health"),
                ("🛠️ الهندسة", "spec_category:engineering"),
                ("💻 علوم الحاسوب", "spec_category:computer"),
                ("📊 العلوم الإدارية", "spec_category:business"),
                ("🗣️ اللغات", "spec_category:languages"),
                ("📿 العلوم الإسلامية", "spec_category:islamic"),
                ("⚖️ القانون والسياسة", "spec_category:law"),
                ("🔬 العلوم الطبيعية", "spec_category:science"),
                ("🎨 الفنون", "spec_category:arts"),
                ("👩‍🏫 التربية", "spec_category:education"),
                ("❓ تخصص آخر", "major_custom"),  # تخصص يدوي
            ]

        # بناء الصفوف بشكل آمن
            for i in range(0, len(main_categories), 2):
                row = main_categories[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            # زر رجوع
            keyboard.add(InlineKeyboardButton("⬅️ رجوع", callback_data="go_account_settings"))

            try:
                bot.edit_message_text(
                    "🎓 *اختر مجال تخصصك*\n\nحدد الفئة الرئيسية أولاً ثم اختر تخصصك الدقيق:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("خطأ في عرض change_specialty:", e)
                try:
                    bot.send_message(chat_id, "🎓 اختر مجال تخصصك:", reply_markup=keyboard)
                except Exception as e2:
                    print("Failed to send change_specialty:", e2)


    # ----------------- عرض التخصصات الفرعية بناءً على الفئة -----------------
        if data.startswith("spec_category:"):
            # صيغة البيانات: spec_category:<key>
            try:
                _, cat_key = data.split(":", 1)
            except ValueError:
                print("spec_category callback غير متوافق:", data)
                return

            SUBS = {  
                "health": [  
                    ("🧬 الطب البشري", "change_major:طب_بشري"),  
                    ("💊 الصيدلة", "change_major:صيدلة"),  
                    ("🏥 التمريض", "change_major:تمريض"),  
                    ("🔬 علوم المختبرات", "change_major:علوم_مختبرات"),  
                    ("🦷 طب الأسنان", "change_major:طب_أسنان"),  
                ],  
                "engineering": [  
                    ("⚙️ الميكانيكا", "change_major:هندسة_ميكانيكية"),  
                    ("🧪 الكيميائية", "change_major:هندسة_كيميائية"),
                    ("💡 الكهربائية", "change_major:هندسة_كهربائية"),  
                    ("🏗️ المدنية", "change_major:هندسة_مدنية"),  
                    ("🔧 البرمجيات", "change_major:هندسة_برمجيات"),  
                    ("📡 الاتصالات", "change_major:هندسة_اتصالات"),  
                ],  
                "computer": [  
                    ("💻 علوم الحاسوب", "change_major:علوم_حاسوب"),  
                    ("📱 تطوير البرمجيات", "change_major:تطوير_برمجيات"),  
                    ("🔒 أمن المعلومات", "change_major:امن_معلومات"),  
                    ("🤖 الذكاء الاصطناعي", "change_major:ذكاء_اصطناعي"),  
                    ("📊 علم البيانات", "change_major:علم_البيانات"),  
                ],  
                "business": [  
                    ("📈 إدارة الأعمال", "change_major:ادارة_اعمال"),  
                    ("💹 الاقتصاد", "change_major:اقتصاد"),  
                    ("📊 المحاسبة", "change_major:محاسبة"),  
                    ("🧮 التسويق", "change_major:تسويق"),  
                    ("🏦 التمويل", "change_major:تمويل"),  
                ],  
                "languages": [  
                    ("🌎 الإنجليزية", "change_major:لغة_انجليزية"),  
                    ("🇫🇷 الفرنسية", "change_major:لغة_فرنسية"),  
                    ("🇸🇦 العربية", "change_major:لغة_عربية"),  
                    ("📚 الترجمة", "change_major:ترجمة"),  
                    ("🇩🇪 الألمانية", "change_major:لغة_ألمانية"),  
                ],  
                "islamic": [  
                    ("📜 الفقه", "change_major:فقه"),  
                    ("💡 العقيدة", "change_major:عقيدة"),  
                    ("📖 التفسير", "change_major:تفسير"),  
                    ("🕌 الدراسات الإسلامية", "change_major:دراسات_اسلامية"),  
                    ("🌙 السيرة النبوية", "change_major:سيرة_نبوية"),  
                ],  
                "law": [  
                    ("📜 القانون الدولي", "change_major:قانون_دولي"),  
                    ("🏛️ القانون الوطني", "change_major:قانون_وطني"),  
                    ("🗳️ العلوم السياسية", "change_major:علوم_سياسية"),  
                    ("👮‍♂️ القانون الجنائي", "change_major:قانون_جنائي"),  
                    ("⚖️ القانون التجاري", "change_major:قانون_تجاري"),  
                ],  
                "science": [  
                    ("🧪 الفيزياء", "change_major:فيزياء"),  
                    ("🔬 الكيمياء", "change_major:كيمياء"),  
                    ("🔢 الرياضيات", "change_major:رياضيات"),  
                    ("🌿 الأحياء", "change_major:احياء"),  
                ],  
                "arts": [  
                    ("🎭 الفنون الأدائية", "change_major:فنون_ادائية"),  
                    ("🖼️ الفنون البصرية", "change_major:فنون_بصرية"),  
                    ("📚 الأدب", "change_major:ادب"),  
                    ("🌏 التاريخ", "change_major:تاريخ"),  
                    ("🎵 الموسيقى", "change_major:موسيقى"),  
                ],  
                "education": [  
                    ("👩‍🏫 علم النفس التربوي", "change_major:علم_النفس_تربوي"),  
                    ("📘 المناهج", "change_major:مناهج_تدريس"),  
                    ("🧩 التربية الخاصة", "change_major:تربية_خاصة"),  
                    ("📊 إدارة التعليم", "change_major:ادارة_تعليم"),  
                    ("👶 الطفولة المبكرة", "change_major:طفولة_مبكرة"),  
                ],  
            }
            sub_btns = SUBS.get(cat_key, [])
            keyboard = InlineKeyboardMarkup(row_width=2)

            # بناء الأزرار بصفين
            for i in range(0, len(sub_btns), 2):
                row = sub_btns[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            # اضافة خيارات النهاية
            keyboard.add(InlineKeyboardButton("❓ تخصص آخر", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("🔙 العودة", callback_data="change_specialty"))
            keyboard.add(InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home"))

            try:
                pretty = cat_key.replace("_", " ")
                bot.edit_message_text(
                f"🎓 *{pretty}*\n\n👇 اختر تخصصك الدقيق من القائمة:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("خطأ في عرض التخصصات الفرعية:", e)
                try:
                    bot.send_message(chat_id, f"🎓 {pretty}\nاختر تخصصك:", reply_markup=keyboard)
                except Exception as e2:
                    print("فشل إرسال فرعية:", e2)


        if data.startswith("change_major:"):
            user_states[uid] = "awaiting_custom_major"
            selected = data.split(":", 1)[1]  # مثال 'طب_بشري'
            save_user_major(uid, selected)     # نفس الدالة لحفظ التخصص في قاعدة البيانات

            bot.edit_message_text(
                f"✅ تم تغيير تخصصك إلى: *{selected.replace('_', ' ')}*",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )

            try:
                send_main_menu(chat_id, message_id)  # إظهار القائمة الرئيسية
            except:
                pass

            return




        if data == "advanced":
            with state_lock:
                user_states[int(uid)] = "awaiting_advanced_test_file"
            bot.answer_callback_query(c.id, "✅ تم تفعيل وضع advanced  — أرسل الملف الآن.")
            bot.edit_message_text(
                "🧠 تم اختيار *الوضع المتقدم*.\n\n"
                "الرجاء إرسال الملف أو النص الذي ترغب بالاختبار منه.",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return

        elif data == "simple":
            with state_lock:
                user_states[uid] = "awaiting_simple_test_file"
            bot.answer_callback_query(c.id, "✅ تم تفعيل وضع simple — أرسل الملف الآن.")
            bot.edit_message_text(
            "📄 تم اختيار *الوضع المباشر*.\n\n"
                "أرسل الملف أو النص للبدء في توليد الاختبار.",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return



        elif data == "game_private":
            try:
                row = fetch_user_major(uid)
                major = row if row else "عام"

                keyboard = InlineKeyboardMarkup(row_width=1)
                keyboard.add(
                    InlineKeyboardButton("🧩 Vocabulary Match", callback_data="game_vocab"),
                    InlineKeyboardButton("⏱️ تحدي السرعة", callback_data="game_speed"),
                    InlineKeyboardButton("❌ الأخطاء الشائعة", callback_data="game_mistakes"),
                    InlineKeyboardButton("🧠 لعبة الاستنتاج", callback_data="game_inference"),
                    InlineKeyboardButton("⬅️ رجوع", callback_data="go_games")
                )
                bot.edit_message_text(
                    f"🎓 تخصصك الحالي: {major}\n"
                    "اختر لعبة 👇",
                    chat_id=chat_id,
                    message_id=message_id,
                    reply_markup=keyboard
                )
            except Exception as e:
                logging.exception("❌ حدث خطأ في game_private")
                bot.send_message(uid, "❌ حدث خطأ أثناء عرض الألعاب.")

    
        elif data == "back_to_games":
            try:
                bot.delete_message(c.message.chat.id, c.message.message_id)
            except Exception as e:
                logging.warning(f"❌ فشل حذف الرسالة عند الرجوع: {e}")
    
    

        elif data in ["game_vocab", "game_speed", "game_mistakes", "game_inference"]:
            game_type = data.split("_", 1)[1]

            # التحقق من إمكانية اللعب اليومي (6 مرات)
            state = game_states.get(uid, {"count": 0})
            if state["count"] >= 6:
                return bot.send_message(uid, "🛑 لقد وصلت إلى الحد الأقصى للألعاب المجانية (6 مرات).")

            if not can_play_game_today(uid, game_type):
                bot.answer_callback_query(c.id, "❌ لقد لعبت هذه اللعبة اليوم!")
                return

            loading_msg = bot.send_message(chat_id, "⏳ جاري تحضير السؤال...")

            try:
                record_game_attempt(uid, game_type)

                # التخصص
                
                row = fetch_user_major(uid)
                major = row if row else "عام"

                # توليد السؤال حسب نوع اللعبة
                if game_type == "vocab":
                    raw = generate_vocabulary_game(uid, major, native_lang="Arabic")
                elif game_type == "speed":
                    raw = generate_speed_challenge(uid, major, native_lang="Arabic")
                elif game_type == "mistakes":
                    raw = generate_common_mistakes_game(uid, major, native_lang="Arabic")
                elif game_type == "inference":
                    raw = generate_inference_game(uid, major, native_lang="Arabic")

                question = raw["question"]
                options = raw["options"]
                correct_index = raw["correct_index"]

                if not isinstance(options, list) or len(options) < 2:
                    raise ValueError("عدد الخيارات غير صالح")

                # حفظ خيارات السؤال في الذاكرة المؤقتة
                game_states[uid] = {"count": state["count"] + 1, "options": options}

                keyboard = InlineKeyboardMarkup(row_width=2)

                # أزرار الإجابات
                for i, option in enumerate(options):
                    short_option = (option[:50] + "...") if len(option) > 50 else option
                    callback_data = f"ans_{game_type}_{i}_{correct_index}"
                    keyboard.add(InlineKeyboardButton(short_option, callback_data=callback_data))
    
                # أزرار التحكم
                keyboard.row(
                    InlineKeyboardButton("🔄 سؤال جديد", callback_data=f"new_{game_type}"),
                    InlineKeyboardButton("⬅️ رجوع", callback_data="back_to_games")
                )
                keyboard.add(
                    InlineKeyboardButton(
                        "📤 شارك هذه اللعبة", 
                        switch_inline_query="جرب هذه اللعبة الرائعة من @Oiuhelper_bot 🎯")
                )

                bot.delete_message(chat_id, loading_msg.message_id)
                text = f"🧠 اختر الإجابة الصحيحة:\n\n{question}"
                bot.send_message(chat_id, text, reply_markup=keyboard)

            except Exception as e:
                try:
                    bot.delete_message(chat_id, loading_msg.message_id)
                except:
                    pass
                logging.error(f"فشل توليد اللعبة: {str(e)}")
                bot.send_message(uid, "❌ حدث خطأ أثناء توليد اللعبة، حاول لاحقاً")

        # معالجة طلب سؤال جديد
    


        elif data.startswith("new_"):
            game_type = data.split("_", 1)[1]

            # تحقق من عدد المحاولات (كما في القسم الرئيسي)
            state = game_states.get(uid, {"count": 0})
            if state["count"] >= 6:
                msg = random.choice([
                    "🚫 وصلت إلى الحد الأقصى لعدد الأسئلة اليوم!\n✨ جرب غدًا أو شارك البوت مع أصدقائك!",
                    "❌ انتهت محاولات اليوم! يمكنك المحاولة مجددًا لاحقًا.",
                    "🛑 لا مزيد من الأسئلة الآن. عد لاحقًا لتكمل رحلتك!"
                ])
                return bot.answer_callback_query(c.id, msg, show_alert=True)

            loading_msg = bot.send_message(c.message.chat.id, "⏳ جاري تحضير السؤال التالي...")

            try:
                # توليد السؤال الجديد
                
                row = fetch_user_major(uid)
                major = row if row else "عام"

                game_generators = {
                    "vocab": generate_vocabulary_game,
                    "speed": generate_speed_challenge,
                    "mistakes": generate_common_mistakes_game,
                    "inference": generate_inference_game
                }

                raw = game_generators[game_type](uid, major)
                question = raw["question"]
                options = raw["options"]
                correct_index = raw["correct_index"]

                if not isinstance(options, list) or len(options) < 2:
                    raise ValueError("عدد الخيارات غير صالح")

                # حفظ خيارات السؤال الجديد
                game_states[uid]["count"] += 1
                game_states[uid]["options"] = options

                # إنشاء الأزرار
                keyboard = InlineKeyboardMarkup(row_width=2)
                for i, option in enumerate(options):
                    short_option = (option[:50] + "...") if len(option) > 50 else option
                    callback_data = f"ans_{game_type}_{i}_{correct_index}"
                    keyboard.add(InlineKeyboardButton(short_option, callback_data=callback_data))

                keyboard.row(
                    InlineKeyboardButton("🔄 سؤال جديد", callback_data=f"new_{game_type}"),
                    InlineKeyboardButton("⬅️ رجوع", callback_data="back_to_games")
                )
                keyboard.add(
                    InlineKeyboardButton(
                        "📤 شارك هذه اللعبة", 
                        switch_inline_query="جرب هذه اللعبة الرائعة من @Oiuhelper_bot 🎯")
                )

                # تعديل نفس الرسالة
                bot.edit_message_text(
                    text=f"🧠 اختر الإجابة الصحيحة:\n\n{question}",
                    chat_id=c.message.chat.id,
                    message_id=c.message.message_id,
                    reply_markup=keyboard
                )

            except Exception as e:
                logging.error(f"❌ فشل توليد سؤال جديد: {e}")
                bot.answer_callback_query(c.id, "❌ فشل توليد السؤال")

            finally:
                try:
                    bot.delete_message(c.message.chat.id, loading_msg.message_id)
                except:
                    pass

        elif data.startswith("ans_"):
            parts = data.split("_")
            game_type = parts[1]
            selected = int(parts[2])
            correct = int(parts[3])

            options = game_states.get(uid, {}).get("options", [])
            correct_text = options[correct] if correct < len(options) else f"الخيار رقم {correct+1}"

            wrong_responses = [
                "❌ خطأ! جرب مجددًا 😉\n✅ الصحيح: {correct}",
                "🚫 للأسف، ليست الصحيحة!\n✅ الجواب: {correct}",
                "😅 ليست الإجابة الصحيحة، الجواب هو: {correct}",
                "❌ لا، حاول مرة أخرى!\n✔️ الصحيح هو: {correct}"
            ]

            if selected == correct:
                bot.answer_callback_query(c.id, "✅ إجابة صحيحة!", show_alert=False)
            else:
                msg = random.choice(wrong_responses).format(correct=correct_text)
                bot.answer_callback_query(c.id, msg, show_alert=False)


        
        # يمكنك إضافة المزيد من المعالجات الأخرى لـ callback_data هنا


        elif data.startswith("soon_"):
            feature_name = {
                "soon_review": "📚 ميزة المراجعة السريعة",
                "soon_summary": "📄 ملخصات PDF",
            }.get(data, "هذه الميزة")

            bot.answer_callback_query(c.id)
            bot.send_message(chat_id, f"{feature_name} ستكون متاحة قريبًا... 🚧")
        
        elif data.startswith("retry:"):
            quiz_code = data[6:]
            quiz_manager.start_quiz(chat_id, quiz_code, bot)
        
        elif data.startswith("share_quiz:"):
            quiz_code = data.split(":", 1)[1]
            chat_id = c.message.chat.id  # ← تأكد من تعيين chat_id هنا

            share_link = f"https://t.me/Oiuhelper_bot?start=quiz_{quiz_code}"

            msg_text = f"""<b>🎉 شارك هذا الاختبار مع زملائك!</b>

        انسخ الرابط أدناه أو اضغط لفتحه مباشرة:
        🔗 <a href="{share_link}">{share_link}</a>

        📝 عند فتح الرابط، سيبدأ الاختبار تلقائيًا بإذن الله.
        """

            keyboard = types.InlineKeyboardMarkup()
            keyboard.add(
                types.InlineKeyboardButton("🔗 نسخ الرابط", switch_inline_query=share_link),
                types.InlineKeyboardButton("🏠 القائمة الرئيسية", callback_data="go_back_home")
            )

            bot.send_message(chat_id, msg_text, parse_mode="HTML", reply_markup=keyboard)

    except Exception as e:
        logging.exception("Callback handler error: %s", e)



# -------------------------------------------------------------------
# ------ message & document Handlers ---------------------------------
# ---------------------------------------


@bot.message_handler(func=lambda m: user_states.get(m.from_user.id) == "awaiting_major", content_types=['text'])
def set_custom_major(msg):
    try:
        major = msg.text.strip()
        uid = msg.from_user.id

        # حفظ التخصص في DB
        save_user_major(uid, major)
    
        # إخطار الأدمن
        bot.send_message(
            ADMIN_ID,
            f"🆕 تخصص جديد أُرسل من المستخدم:\n"
            f"👤 @{msg.from_user.username or msg.from_user.id}\n"
            f"📚 التخصص: {major}"
        )
        
        # تغيير الحالة للانتظار لاحقاً
        user_states[uid] = "awaiting_simple_test_file"
    
        # إرسال رسالة تأكيد
        bot.send_message(
            uid,
            f"✅ تم تحديد تخصصك: {major}\n"
            f"الآن أرسل ملف (PDF/DOCX/TXT) أو نصًا مباشرًا لتوليد اختبارك."
        )
    except Exception as e:
        logging.error(f"Error in set_custom_major: {e}")
        bot.send_message(uid, "❌ حدث خطأ أثناء معالجة طلبك، يرجى المحاولة لاحقاً")
@bot.message_handler(func=lambda m: user_states.get(m.from_user.id) in [
    "awaiting_major_for_games",
    "awaiting_custom_major"
])
def handle_user_major(msg):
    if msg.chat.type != "private":
        return  # تجاهل الرسائل في المجموعات

    uid = msg.from_user.id
    state = user_states.get(uid)
    major = msg.text.strip()
    if uid in user_states and 'message_id' in user_states[uid]:
        stored_message_id = user_states[uid]['message_id']

    if len(major) < 2:
        bot.send_message(uid, "⚠️ يرجى إدخال تخصص صالح.")
        return

    cursor.execute("INSERT OR REPLACE INTO users(user_id, major) VALUES(?, ?)", (uid, major))
    conn.commit()
    user_states.pop(uid, None)


    if state == "awaiting_major_for_games":
        bot.send_message(uid, f"✅ تم تسجيل تخصصك: {major}\n"
                              "الآن يمكنك اختيار لعبة من قائمة الألعاب التعليمية.")
        keyboard = InlineKeyboardMarkup(row_width=1)
        keyboard.add(
            InlineKeyboardButton("🔒 العب في الخاص", callback_data="game_private"),
            InlineKeyboardButton("👥 العب في المجموعة", switch_inline_query="game")
        )
        bot.send_message(uid, "🎮 اختر طريقة اللعب:", reply_markup=keyboard)

    elif state == "awaiting_custom_major":
        sent = bot.send_message(uid, f"✅ تم تسجيل تخصصك: *{major}*", parse_mode="Markdown")
        time.sleep(2)
        try:
            bot.edit_message_text(
                "⬇️ هذه هي القائمة الرئيسية:",
                chat_id=sent.chat.id,
                message_id=sent.message_id
            )
            send_main_menu(uid, message_id=sent.message_id)
        except:
            send_main_menu(uid)



# مجرد إقرار سريع في الـ handler ثم وضع في الطابور

@bot.message_handler(content_types=['text', 'document', 'photo'])
def unified_handler(msg):
    if msg.chat.type != "private":
        return
    
    uid = int(msg.from_user.id)

    file_id = getattr(getattr(msg, 'document', None), 'file_id', None)
    if is_request_already_queued(file_id=file_id, user_id=uid, message_id=msg.message_id):
        # سجل/أرسل للمستخدم رسالة قصيرة تفيد أن الطلب قيد الانتظار بالفعل
        bot.reply_to(msg, "⏳ طلبك قيد المعالجة بالفعل.")
        return



    # حفظ الطلب إن أردت
    try:
        save_request(msg)
    except Exception:
        pass
    
    update_files_and_users(uid, files_count=1)
    update_daily_stats(files=1)


    sent_msg = bot.reply_to(msg, "📝 جاري توليد المحتوى، يرجى الانتظار قليلاً...")
    try:

        request_queue.put_nowait((msg, sent_msg))
    except queue.Full:
        bot.edit_message_text(" العديد من الأشخاص يقومون بتوليد محتوى في الوقت الحالي، يرجى الانتظار قليلاً ️🕰️. الرجاء الإنتظار ...", chat_id=sent_msg.chat.id, message_id=sent_msg.message_id)
        time.sleep(random.randint(1, 2))
        while True:
            try:
                request_queue.put((msg, sent_msg), timeout=5)
            
                break
            except Exception as e:
                bot.edit_message_text("حدث خطأ، يرجى الانتظار قليلاً...", chat_id=sent_msg.chat.id, message_id=sent_msg.message_id)
                time.sleep(1)


# -------------------------------------------------------------------
# الدالة التي تنفذ فعليًا المعالجة (تعمل داخل العامل)
# -------------------------------------------------------------------
def process_message(msg, message_id=None, chat_id=None):
    logging.info("process_message enter: uid=%s type=%s", msg.from_user.id, msg.content_type)

    # ...

    content_type = msg.content_type
    username = msg.from_user.username or "بدون اسم مستخدم"
    uid = msg.from_user.id
        
    with semaphore:
        if msg.content_type == "document":
            file_id = msg.document.file_id
        else:
            file_id = None  # لا تحدّث requests هنا لأن لا ملف


        try:
            if file_id:
                update_request_status(file_id, 'processing')
                print(f"🚀 بدء معالجة: {file_id}")
            time.sleep(random.randint(1, 2))  # محاكاة معالجة
            
            if file_id:
                update_request_status(file_id, 'done')
                print(f"✅ تم المعالجة: {file_id}")
        except Exception:
            logging.exception("requests status update failed")

        


    # اجلب الحالة
    try:
        with state_lock:
            state = user_states.get(int(uid))
    except NameError:
        # إذا ما عندك state_lock، استخدم بدون قفل
        state = user_states.get(int(uid))

    logging.info("state for uid %s: %s", uid, state)

    major = fetch_user_major(uid)
    

    content = ""
    path = None

    # log at entry
    logging.info("Message received: uid=%s type=%s text_len=%s", uid, content_type, len(getattr(msg, 'text', '') or ""))


    
        # فور التأكد من وجود حالة، اعرض للّوج
    if state is None:
        # يمكن الرد للتجربة فقط أثناء الديباغ
        # bot.send_message(uid, f"DEBUG: no state set (you are {uid})")
        return


    try:
        if content_type == "text":
            content = msg.text or ""
            coverage = "كاملة ✅"

        
        # معالجة الصور (photo)
        elif msg.content_type == "photo":
            if not can_generate(uid):
                return bot.reply_to(msg, "⚠️ هذه الميزة متاحة فقط للمشتركين.")
            
            file_id = msg.photo[-1].file_id
            file_info = bot.get_file(file_id)
            file_data = bot.download_file(file_info.file_path)

            os.makedirs("downloads", exist_ok=True)
            path = os.path.join("downloads", f"{uid}_photo.jpg")
            with open(path, "wb") as f:
                f.write(file_data)

            bot.edit_message_text("🖼️ جاري استخراج النص من الصورة...", chat_id=chat_id, message_id=message_id)


            content, ocr_debug = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng")
            if not content.strip():
                return bot.send_message(uid, f"❌ فشل في استخراج النص من الصورة. {ocr_debug}")


        elif msg.content_type == "document":
            # وضع الطلب في القائمة ليتم معالجته بواسطة العمال
            file_info = bot.get_file(msg.document.file_id)
            if file_info.file_size > 5 * 1024 * 1024:
                return bot.send_message(uid, "❌ الملف كبير جدًا، الحد 5 ميغابايت.")
    
            file_data = bot.download_file(file_info.file_path)
            os.makedirs("downloads", exist_ok=True)
            path = os.path.join("downloads", msg.document.file_name)

            with open(path, "wb") as f:
                f.write(file_data)
            user_files[uid] = path

            ext = path.rsplit(".", 1)[-1].lower()
            # بعد الاستخراج العادي
            if ext == "pdf":
                content_full = extract_text_from_pdf(path)  # النص الكامل
                full_length = len(content_full)

                # إذا المستخدم غير مشترك، اقتطع فقط 3000 حرف
                if not can_generate(uid):
                    content = content_full[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% من الملف"
                else:
                    content = content_full
                    coverage = "كاملة ✅"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(
                            uid,
                            "⚠️ لا يمكن قراءة هذا الملف تلقائيًا. تتطلب المعالجة المتقدمة اشتراكًا فعالًا."
                        )
                    bot.reply_to(msg, "⏳ يتم تجهيز الملف... الرجاء الانتظار لحظات.")
                    language = detect_language_from_filename(msg.document.file_name)
                    content, ocr_debug = extract_text_from_pdf_with_ocr(path, api_key=OCR_API_KEY, language=language)
                    if not content.strip():
                        bot.send_message(uid, f"❌ فشل في استخراج النص من الملف. {ocr_debug}")
                        return
                    preview = content[:1500]
                    bot.send_message(uid, f"📄 تم استخراج النص بنجاح (جزء منه):\n\n{preview}")
            elif ext == "docx":
                content_full = extract_text_from_pdf(path)  # النص الكامل
                full_length = len(content_full)
                # إذا المستخدم غير مشترك، اقتطع فقط 3000 حرف
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% من الملف"
                else:
                    content = content_full
                    coverage = "كاملة ✅"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "⚠️ لا يمكن قراءة هذا الملف تلقائيًا. تتطلب المعالجة المتقدمة اشتراكًا فعالًا.")
                    bot.edit_message_text("⏳ يتم تجهيز الملف... الرجاء الانتظار لحظات.", chat_id=chat_id, message_id=message_id)
                    language = detect_language_from_filename(msg.document.file_name)
                    content = extract_text_from_pdf_with_ocr(path, api_key=OCR_API_KEY, language=language)
            elif ext == "txt":
                content_full = extract_text_from_pdf(path)  # النص الكامل
                full_length = len(content_full)
                # إذا المستخدم غير مشترك، اقتطع فقط 3000 حرف
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% من الملف"
                else:
                    content = content_full
                    coverage = "كاملة ✅"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "⚠️ لا يمكن قراءة هذا الملف تلقائيًا. تتطلب المعالجة المتقدمة اشتراكًا فعالًا.")
                    bot.edit_message_text("⏳ يتم تجهيز الملف... الرجاء الانتظار لحظات.", chat_id=chat_id, message_id=message_id)
                    content = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng+ara")
                
            elif ext == "pptx":
                content_full = extract_text_from_pdf(path)  # النص الكامل
                full_length = len(content_full)
                
                # إذا المستخدم غير مشترك، اقتطع فقط 3000 حرف
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% من الملف"
                else:
                    content = content_full
                    coverage = "كاملة ✅"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "⚠️ لا يمكن قراءة هذا الملف تلقائيًا. تتطلب المعالجة المتقدمة اشتراكًا فعالًا.")
                    bot.edit_message_text("⏳ يتم تجهيز الملف... الرجاء الانتظار لحظات.", chat_id=chat_id, message_id=message_id)
                    language = detect_language_from_filename(msg.document.file_name)
                    content = extract_text_from_pptx_with_ocr(path, api_key=OCR_API_KEY, language=language)

            elif ext in ("jpg", "png"):
                if not can_generate(uid):
                    return bot.send_message(uid, "⚠️ هذه الميزة متاحة فقط للمشتركين.")
                bot.edit_message_text("⏳ جاري تحليل الصورة...", chat_id=chat_id, message_id=message_id)
                content, ocr_debug = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng")
            
               
            else:
                return bot.send_message(uid, "⚠️ نوع الملف غير مدعوم. أرسل PDF أو Word أو TXT.")

        else:
            try:
                os.remove(path)
            except Exception as e:
                print(f"[WARNING] لم يتم حذف الملف المؤقت: {e}")
        if not state:
        # لا حالة: لا تفعل شيء
            return

        if not content or not content.strip():
            return bot.send_message(uid, "⚠️ لم أتمكن من قراءة محتوى الملف أو النص.")
        print(f">>> Content preview: {content[:300]}")

        waiting_messages_anki = [
            "🧠 يتم تحليل النص لاستخلاص النقاط الأساسية...",
            "✨ جاري تحويل المحتوى إلى بطاقات ذكية سهلة المراجعة...",
            "📚 تنظيم المعلومات لتعزيز قدرتك على التذكر...",
            "🎨 تصميم البطاقات بأسلوب يساعد على التعلم السريع...",
            "⚙️ إعداد البطاقات لتكون جاهزة للمراجعة الفعالة...",
            "🔍 التركيز على المفاهيم الجوهرية لبناء أساس قوي...",
            "🚀 تحسين هيكل البطاقات لتجربة دراسية سلسة...",
            "💡 جاري صياغة الأسئلة والأجوبة بوضوح ودقة..."
        ]


        waiting_messages_quiz = [
            "📝 جاري صياغة أسئلة دقيقة من المحتوى المقدم...",
            "🎯 تحديد أهم الأفكار لوضع أسئلة تقيس الفهم الحقيقي...",
            "⚖️ موازنة أنواع الأسئلة لتغطية شاملة للموضوع...",
            "🧠 تصميم اختبار ذكي يتحدى معرفتك بشكل بنّاء...",
            "📊 تنظيم الأسئلة لتقديم تقييم متكامل لمستواك...",
            "🧩 بناء اختبار متماسك يقيس مختلف جوانب المعرفة...",
            "🔍 فحص المحتوى بعناية لصياغة أسئلة محكمة...",
            "✅ التحقق من جودة الأسئلة لضمان اختبار عادل..."
        ]


        progress_messages = [
            "⏳ لحظات ونبدأ... 25% من الطريق",
            "⚡️ المعالجة في منتصفها... 50% من الطريق",
            "🚀 اقتربنا من الإنجاز... 75% من الطريق",
            "🎉 اللمسات الأخيرة... 90% من الطريق"
        ]




         # ============================
        # Awaiting AI Anki
        # ============================

        if state == "awaiting_anki_file_ai":
            logging.info("Handling awaiting_anki_file_ai for uid=%s", uid)
            if not can_generate(uid):
                return bot.send_message(uid, "⚠️ لقد استنفدت 3 اختبارات مجانية هذا الشهر.")


    
            # إعداد رسالة التحميل الأولية
            loading_msg = safe_edit_or_send("🔄 جاري معالجة الملف...", chat_id, message_id)
    
            try:
                # إذا كان المحتوى كبيراً
                if len(content) > 10000:
                    try:
                        # تحديث الرسالة لعملية التلخيص
                        bot.edit_message_text(
                            chat_id=chat_id,
                            message_id=message_id,
                            text="📚 المحتوى كبير جداً\n🔍 جاري تلخيص المحتوى..."
                        )
                
                        content = summarize_long_text(content)
                
                        # تأكيد نجاح التلخيص
                        bot.edit_message_text(
                            chat_id=uid,
                            message_id=loading_msg.message_id,
                            text="✅ تم تلخيص المحتوى بنجاح\n⏳ جاري إنشاء البطاقات..."
                        )
                        time.sleep(1)
                
                    except Exception as e:
                        print("[ERROR] فشل في تلخيص المحتوى:", e)
                        return bot.edit_message_text(
                            chat_id=uid,
                            message_id=loading_msg.message_id,
                            text="❌ فشل في تلخيص المحتوى.\n\nيرجى إرسال ملف أصغر أو المحاولة لاحقاً."
                        )
        
                # مؤشر تقدم متحرك
                progress_phrases = [
                    "📖 جاري تحليل المحتوى...",
                    "🧠 معالجة المعلومات...",
                    "🛠️ إنشاء البطاقات...",
                    "✨ جاري التنسيق النهائي..."
                ]
        
                for i, phrase in enumerate(progress_phrases):
                    # إضافة شريط تقدم بصري
                    progress_bar = "[" + "=" * (i+1) + " " * (len(progress_phrases)-i-1) + "]"
            
                    bot.edit_message_text(
                        chat_id=uid,
                        message_id=loading_msg.message_id,
                        text=f"{progress_bar}\n\n{phrase}\n\n⏳ يرجى الانتظار..."
                    )
                    time.sleep(1.5)
        
                # إضافة رسالة انتظار جذابة
                bot.edit_message_text(
                    chat_id=uid,
                    message_id=loading_msg.message_id,
                    text=f"🎯 {random.choice(waiting_messages_anki)}\n\n⚡ جاري الانتهاء من التحضير..."
                )
                time.sleep(random.randint(2, 5))
        
                # إنشاء البطاقات
                if not can_generate(uid):
                    cards, title = generate_anki_cards_from_text(content, major=major, user_id=uid)
                    
                else:
                    cards, title = generate_special_anki_cards_from_text(content, major=major, user_id=uid)

                if not cards:
                    return bot.edit_message_text(
                        chat_id=uid,
                        message_id=loading_msg.message_id,
                        text="❌ لم أتمكن من إنشاء أي بطاقات.\n\nقد يكون المحتوى غير مناسب أو حدث خطأ أثناء المعالجة."
                    )
                    
        
                # تنظيف العنوان ليكون اسم ملف صالح
                safe_title = re.sub(r'[^a-zA-Z0-9_\u0600-\u06FF]', '_', title)[:40]
                filename = f"{title}_{timestamp}.apkg".replace(" ", "_")
        
                # حفظ الملف وإرساله مع تحديث الرسالة السابقة
                filepath = save_cards_to_apkg(cards, filename=filename, deck_name=title)
        
                 # تحرير الرسالة الأخيرة لإظهار نجاح العملية
                bot.edit_message_text(
                    chat_id=uid,
                    message_id=loading_msg.message_id,
                    text=f"✅ تم إنشاء {len(cards)} بطاقة بنجاح!\n\n📚 العنوان: {title}\n\n⚡ جاري إرسال الملف..."
                )
                increment_count(uid)
                notify_admin("توليد أنكي آلي", username, uid)

                # إرسال الملف مع caption
                with open(filepath, 'rb') as file:
                    bot.send_document(
                        chat_id=uid,
                        document=file,
                        caption=f"📂 {title}\n\n🎴 عدد البطاقات: {len(cards)}\n\nاستمتع بالدراسة!",
                        reply_to_message_id=loading_msg.message_id
                    )
                    with state_lock:
                        user_states.pop(uid, None)
                    logging.info("Finished ai_anki for uid=%s", uid)



            except Exception:
                logging.exception("Error while processing ai anki for uid=%s", uid)
                with state_lock:
                    user_states.pop(uid, None)
                    bot.send_message(uid, "حدث خطأ أثناء إنشاء البطاقات. حاول لاحقًا.")
            return


        # ============================
        # Awaiting manual anki
        # ============================
        elif state == "awaiting_anki_file_manual":
            if msg.content_type == "text":
            # قائمة رسائل انتظار متحركة
                waiting_messages = [
                    "🌱 جارٍ تجهيز مدخلاتك التعليمية...",
                    "🧠 النظام المعرفي يعمل بكامل طاقته...",
                    "🔮 يتم تنظيم المحتوى التعليمي...",
                    "🚀 يجري الإعداد لمرحلة الإطلاق...",
                    "🎩 تتم معالجة البيانات بصورة دقيقة..."
                ]
        
                # إرسال رسالة الانتظار مع مؤثرات بصرية
                waiting_msg = bot.send_message(chat_id, "⏳ **جارٍ المعالجة**\n`0% اكتمال`", 
                                      parse_mode="Markdown")
        
                # خطوات التقدم مع رموز إبداعية
                processing_steps = [
                    {"icon": "🔍", "text": "تحليل النصوص المدخلة", "delay": 0.8},
                    {"icon": "🧩", "text": "بناء البنية التعليمية", "delay": 1.2},
                    {"icon": "🎨", "text": "إعداد بطاقات المراجعة", "delay": 1.0},
                    {"icon": "⚡", "text": "تجهيز الحزمة النهائية", "delay": 0.7},
                    {"icon": "🚀", "text": "إطلاق الملف التعليمي", "delay": 1.5}
                ]
        
                # محاكاة التقدم التدريجي
                progress = 0
                step_size = 100 // len(processing_steps)
        
                for idx, step in enumerate(processing_steps):
                    # حساب النسبة المئوية
                    progress = min(100, (idx + 1) * step_size)
                    progress_bar = "🟩" * (progress // 10) + "⬜" * (10 - progress // 10)
            
                    # بناء الرسالة مع تأثير التراكم
                    message_text = (
                        f"⏳ **جارٍ المعالجة**\n"
                        f"`{progress}% اكتمال`\n"
                        f"{progress_bar}\n\n"
                        f"{step['icon']} **المرحلة {idx+1}:** {step['text']}"
                    )
            
                    # تحديث الرسالة مع تأثير التدرج
                    try:
                        bot.edit_message_text(
                            message_text,
                            chat_id=waiting_msg.chat.id,
                            message_id=waiting_msg.message_id,
                            parse_mode="Markdown"
                )
                    except:
                        pass
            
                    # تأخير ديناميكي بين الخطوات
                    time.sleep(step['delay'])
        
                # معالجة الملف الفعلية
                cards = parse_manual_anki_input(msg.text)
                if cards:
                    # إنشاء الملف
                    output_file = f"{uid}_manual_anki.apkg"
                    save_cards_to_apkg(cards, filename=output_file, deck_name="مكتبتك التعليمية")
            
                    # إرسال الملف مع رسالة رسمية
                    with open(output_file, 'rb') as file:
                        bot.send_document(
                            chat_id=uid,
                            document=file,
                            caption=(
                                f"🌿 *تم إنشاء ملفك التعليمي بنجاح.*\n"
                                f"عدد البطاقات: {len(cards)} بطاقة\n"
                                f"مدة التنفيذ: {random.randint(3,7)} ثوانٍ\n\n"
                                f"📚 ملف المراجعة جاهز للاستخدام."
                            ),
                            reply_to_message_id=message_id,
                            parse_mode="Markdown"
                        )
                        notify_admin("توليد أنكي يدوي", username, uid)

            
                    # حذف رسالة التقدم بعد الإرسال
                    try:
                        bot.delete_message(chat_id, waiting_msg.message_id)
                    except:
                        pass
            
                    # رسالة ختامية رسمية
                    if random.random() < 0.1:  # 30% احتمال
                        bot.send_message(
                            chat_id,
                            "✨ *أحسنت! تم إنشاء بطاقاتك بنجاح.*\n"
                            "🚀 واصل المراجعة بانتظام، وستتفاجأ بسرعة تقدمك.\n"
                            "─── ⋆⋅☆⋅⋆ ──\n"
                            "💡 *تذكير:* راجع بطاقاتك غدًا، فالتكرار هو سر تثبيت المعرفة!",
                            reply_to_message_id=message_id,
                            parse_mode="Markdown"
                    )
                    with state_lock:
                        user_states.pop(uid, None)

                else:
                    # حذف رسائل الانتظار والخطوات وإرسال رسالة الخطأ
                    bot.delete_message(uid, waiting_msg.message_id)
                    for step_msg in step_messages:
                        bot.delete_message(uid, step_msg.message_id)
                    bot.send_message(uid, "❌ لم يتم العثور على أي بطاقات صالحة")
            else:
                bot.send_message(uid, "❌ يرجى إرسال النص فقط لإنشاء بطاقات Anki يدويًا.")

        # ============================
        # Awaiting advanced test file
        # ============================
        # استخدم مكتبة traceback لتشخيص دقيق
        elif state == "awaiting_advanced_test_file":
            import traceback
    
            # احتفظ بالـ ID الأصلي للرسالة والدردشة في متغيرات آمنة
            original_chat_id = msg.chat.id
            original_message_id = message_id

            try:
                # 1. ابدأ بتعديل الرسالة الأولى
                bot.edit_message_text("🤖 جاري معالجة الملف وإنشاء اختبار ذكي...", chat_id=original_chat_id, message_id=original_message_id)
        
                if not can_generate(uid):
                    return bot.send_message(uid, "⚠️ لقد استنفدت 3 اختبارات مجانية هذا الشهر.")

                if len(content) > 10000:
                    bot.edit_message_text("🔍 المحتوى كبير، جاري تلخيصه...", chat_id=original_chat_id, message_id=original_message_id)
                    content = summarize_long_text(content)
        
                # 2. عرض رسائل التقدم (لا نغير قيمة المتغيرات الأصلية)
                bot.edit_message_text("🧠 جاري توليد الاختبار، الرجاء الانتظار...", chat_id=original_chat_id, message_id=original_message_id)
                for progress_msg in progress_messages:
                    bot.edit_message_text(progress_msg, chat_id=original_chat_id, message_id=original_message_id)
                    time.sleep(1.5)

                bot.edit_message_text(random.choice(waiting_messages_quiz), chat_id=original_chat_id, message_id=original_message_id)
                time.sleep(2)

                # 3. توليد الاختبار
                print("[ADVANCED_QUIZ] بدء توليد الاختبار الطبي المتقدم...")
                quiz_data = generate_Medical_quizzes(content=content, major="General Medicine", user_id=uid)
        
                # طباعة للتحقق من القيم قبل الإرسال النهائي
                print(f"[DEBUG] chat_id: {original_chat_id}, message_id: {original_message_id}, quiz_data is not None: {quiz_data is not None}")
    
                if quiz_data:
                    # 4. إرسال النتيجة النهائية
                    send_quiz_to_user(original_chat_id, quiz_data)
                    
                    notify_admin("توليد اختبار ذكي", username, uid)
                    update_top_user(uid, tests=1)
                else:
                    bot.edit_message_text("❌ فشل في إنشاء الاختبار. قد يكون المحتوى غير مناسب. يرجى المحاولة لاحقاً.", chat_id=original_chat_id, message_id=original_message_id)

            except Exception as e:
                # طباعة الخطأ الكامل في الكونسول لتشخيصه
                print("!!!!!!!!!!!!!!!!!! خطأ فادح في معالجة الاختبار المتقدم !!!!!!!!!!!!!!!!!!")
                traceback.print_exc()
                print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        
                error_message = "⚠️ حدث خطأ تقني. تم إبلاغ المطورين."
                try:
                    bot.edit_message_text(error_message, chat_id=original_chat_id, message_id=original_message_id)
                except:
                    bot.send_message(original_chat_id, error_message)
            finally:
                # تأكد من إزالة حالة المستخدم لتجنب بقائه عالقاً
                with state_lock:
                    user_states.pop(uid, None)

                
        # ============================
        # Awaiting simple test file
        # ============================
        elif state == "awaiting_simple_test_file":
            print("[QUIZ] بدء توليد الاختبار للمستخدم:", uid)
            if not can_generate(uid):
                return bot.send_message(uid, "⚠️ لقد استنفدت 3 اختبارات مجانية هذا الشهر.")


            if len(content) > 10000:
                loading_msg = bot.edit_message_text("🔍 المحتوى كبير، جاري تلخيصه...", chat_id=chat_id, message_id=message_id)
                try:
                    print("[QUIZ] المحتوى كبير، جاري التلخيص...")
                    content = summarize_long_text(content)
                except Exception as e:
                    print("[ERROR] تلخيص المحتوى فشل:", e)
                    return bot.send_message(uid, "❌ فشل في تلخيص المحتوى. أرسل ملفًا أصغر أو حاول لاحقًا.")

            else:
                loading_msg = bot.edit_message_text("🧠 جاري توليد الاختبار، الرجاء الانتظار...", chat_id=chat_id, message_id=message_id)

                # عرض رسائل التحميل
            for progress_msg in progress_messages:
                try:
                    bot.edit_message_text(chat_id=uid, message_id=loading_msg.message_id, text=progress_msg)
                except Exception as e:
                    logging.exception("[QUIZ] فشل أثناء تحديث رسالة التقدّم")

                    time.sleep(1.5)
            try:
                bot.edit_message_text(chat_id=uid, message_id=loading_msg.message_id,
                              text=random.choice(waiting_messages_quiz))
            except Exception:
                pass

            time.sleep(2)

            print("[QUIZ] استدعاء generate_quizzes_from_text...")
            quizzes = generate_quizzes_from_text(content, major=major, user_id=uid, num_quizzes=10)
            print("[QUIZ] رجع:", type(quizzes), "بطول:", len(quizzes) if quizzes else "None")

        
            if isinstance(quizzes, list) and len(quizzes) > 0:
                
                try:
                    print(f"تم توليد {len(quizzes)} سؤالا")
                    # تخزين الاختبار أولاً
                    quiz_code = store_quiz(uid, quizzes, bot)
                    print("[QUIZ] كود الاختبار:", quiz_code)

                    if not quiz_code:
                        raise Exception("Failed to store quiz")
                        
                    waiting_quiz = loading_msg.message_id
                    major = fetch_user_major(uid)
                    file_path = user_files[uid]
                    level = "متوسط"

                    # إرسال رسالة "إختبارك جاهز" مع رابط الاختبار
                    quiz_link = f"https://t.me/QuizzyAI_bot?start=quiz_{quiz_code}"
                    estimated_time = len(quizzes) * 30

                    # إرسال رسالة "إختبارك جاهز" مع رابط الاختبار
                    markup = InlineKeyboardMarkup()
                    btn = InlineKeyboardButton("فتح الاختبار", url=quiz_link)
                    markup.add(btn)

                    quiz_msg = (
                    "✨✔️ <b>إختبارك جاهز!</b>\n"
                    "──────────────────\n"
                    f"📂 <b>العنوان:</b> {msg.document.file_name}\n\n"
                    f"📋 <b>عدد الأسئلة:</b> {len(quizzes)}\n"
                    f"⏱️ <b>الزمن الكلي:</b> {estimated_time // 60} دقيقة و {estimated_time % 60} ثانية\n"
                    f"🎓 <b>التخصص:</b> {major} \n"
                    "📦 <b>نوع الاختبار:</b> خاص\n\n"
                    f"📉 <b>التغطية:</b> {coverage}\n"
                    "💡 <b>ميزة الشرح:</b> غير متوفرة\n"
                    f"📊 <b>المستوى:</b> {level}\n\n"
                    "❓هل أنت جاهز للإختبار\n"
                    f"👈 <a href=\"{quiz_link}\">اضغط هنا للبدء</a>"
                    )
                    try:
                        bot.delete_message(chat_id=chat_id, message_id=loading_msg.message_id)
                    except Exception as del_err:
                        print(f"لم يتمكن من حذف رسالة التحميل: {del_err}")
                
                    bot.send_message(chat_id, quiz_msg, reply_markup=markup, parse_mode="HTML", disable_web_page_preview=True)
                    

                    update_top_user(uid, tests=1)
                    notify_admin("توليد اختبار", username, uid)
                    
                    

                    with state_lock:
                        user_states.pop(uid, None)

                except Exception as e:
                    print(f"Error in quiz generation: {e}")
                    bot.send_message(uid, "حدث خطأ غير متوقع   .")
                    
            else:
                print("[QUIZ] التوليد فشل أو رجع None")
                bot.send_message(uid, "❌ لم يتمكن البوت من توليد الاختبار.")


        else:
            bot.reply_to(msg, "عذرًا، لا أستطيع فهم طلبك في الوقت الحالي. يرجى المحاولة لاحقًا أو الاتصال بالدعم.")
            # أو يمكنك تسجيل ذلك في قاعدة البيانات أو ملف السجلات
            print("No state for user", uid, "state:", state)


    except Exception as e:
        import traceback
        logging.exception("process_message error: %s", e)
        print("!!!!!!!!!!!!!!!!! حدث خطأ !!!!!!!!!!!!!!!!!!")
        traceback.print_exc() # هذا السطر سيطبع الخطأ الكامل ومكانه بالضبط
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        bot.send_message(uid, "حدث خطأ غير متوقع   .")
    finally:
        # حذف الملف المؤقت إن وُجد
        if path and os.path.exists(path):
            try:
                os.remove(path)
            except Exception as e:
                print(f"[WARNING] لم يتم حذف الملف المؤقت: {e}")



known_channels = set()


@bot.channel_post_handler(func=lambda msg: True)
def handle_channel_post(msg):
    channel_id = msg.chat.id

    if channel_id in known_channels:
        return  # تم معالجته من قبل

    known_channels.add(channel_id)

    try:
        bot.send_message(
            ADMIN_ID,
            f"📢 تم اكتشاف قناة جديدة:\n\n"
            f"*الاسم:* {msg.chat.title}\n"
            f"*ID:* `{channel_id}`",
            parse_mode="Markdown"
        )
    except Exception as e:
        print(f"[ERROR] إرسال المعرف فشل: {e}")



# -------------------------------------------------------------------
#                           Run Bot
# -------------------------------------------------------------------

@app.route('/anki_preview')
def anki_preview():
    user_cards = generate_anki_cards_from_text(text)[:5]  # ← نحصل على أول 5 بطاقات
    session['cards'] = user_cards
    session['index'] = 0
    session['show_back'] = False
    return redirect('/anki')
    
app.secret_key = 'anki_secret'  # سر الجلسة لتخزين البيانات مؤقتًا


@app.route('/anki', methods=['GET', 'POST'])
def anki_cards():
    content = session.get('anki_content')
    major = session.get('anki_major', 'General')
    if 'cards' not in session:
        session['cards'] = example_cards[:5]
        session['index'] = 0
        session['show_back'] = False

    if request.method == 'POST':
        action = request.form.get('action')
        if action == 'show':
            session['show_back'] = True
        elif action == 'next':
            session['index'] += 1
            session['show_back'] = False

    index = session['index']
    cards = session['cards']

    if index >= len(cards):
        session.clear()
        return "<h2>🎉 انتهيت من البطاقات! أحسنت.</h2><a href='/anki'>🔁 ابدأ من جديد</a>"

    return render_template('anki_viewer.html',
                           card=cards[index],
                           index=index,
                           total=len(cards),
                           show_back=session['show_back'])
# بدء البوت
start_workers()



import json
from datetime import datetime

def insert_sample_quiz_if_not_exists(db_path='quiz_users.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()
        
    """
    Inserts a sample quiz into the database if it doesn't already exist.

    Args:
        cursor: The database cursor object.
        conn: The database connection object.
    """
    cursor.execute("SELECT quiz_code FROM sample_quizzes WHERE quiz_code = ?", ("sample",))
    if cursor.fetchone() is None:
        # The list of questions should be a single list of dictionaries, not a list within a list.
        sample_quiz_data = [
            {
                "question": "ما هو أطول برج في العالم؟",
                "options": ["برج خليفة", "برج إيفل", "برج بيزا", "برج شنغهاي"],
                "correct_index": 0,
                "explanation": "برج خليفة في دبي هو أطول برج في العالم منذ اكتماله عام 2010."
            },
            {
                "question": "ما هو مجموع 7 + 5؟",
                "options": ["10", "12", "13", "14"],
                "correct_index": 1,
                "explanation": "7 + 5 = 12."
            },
            {
                "question": "ما هي عاصمة فرنسا؟",
                "options": ["باريس", "روما", "برلين", "مدريد"],
                "correct_index": 0,
                "explanation": "باريس هي عاصمة فرنسا وأشهر مدنها."
            },
            {
                "question": "ما هو أكبر محيط في العالم؟",
                "options": ["المحيط الأطلسي", "المحيط الهندي", "المحيط الهادئ", "المحيط المتجمد الشمالي"],
                "correct_index": 2,
                "explanation": "المحيط الهادئ هو أكبر محيط على الأرض."
            },
            {
                "question": "كم عدد أيام الأسبوع؟",
                "options": ["5", "6", "7", "8"],
                "correct_index": 2,
                "explanation": "الأسبوع يحتوي على 7 أيام."
            },
            {
                "question": "ما هو الكوكب الأحمر؟",
                "options": ["المشتري", "المريخ", "الزهرة", "عطارد"],
                "correct_index": 1,
                "explanation": "المريخ يسمى الكوكب الأحمر بسبب لونه."
            },
            {
                "question": "ما هو الغاز الذي نتنفسه؟",
                "options": ["الأكسجين", "ثاني أكسيد الكربون", "الهيدروجين", "النيتروجين"],
                "correct_index": 0,
                "explanation": "الأكسجين هو الغاز الأساسي الذي نتنفسه."
            },
            {
                "question": "كم عدد الحواس عند الإنسان؟",
                "options": ["4", "5", "6", "7"],
                "correct_index": 1,
                "explanation": "الإنسان لديه خمس حواس رئيسية."
            }
        ]
        
        # Convert the Python list of dictionaries to a JSON string.
        # `ensure_ascii=False` is important for correctly handling Arabic characters.
        sample_quiz_json = json.dumps(sample_quiz_data, ensure_ascii=False)
        
        cursor.execute(
            "INSERT INTO sample_quizzes (quiz_code, quiz_data, created_at) VALUES (?, ?, ?)",
            ("sample", sample_quiz_json, datetime.utcnow().isoformat())
        )
        conn.commit()
    conn.close()




# نقطة نهاية الويب هوك
@app.route('/' + os.getenv('BOT_TOKEN'), methods=['POST'])
def webhook_bot():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot.process_new_updates([update])
        return 'ok', 200
    return 'Method Not Allowed', 405



@app.route('/' + os.getenv('BOT_TOKEN_2'), methods=['POST'])
def webhook_bot2():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot2.process_new_updates([update])
        return 'ok', 200

    return 'Method Not Allowed', 405

@app.route('/' + os.getenv('BOT_TOKEN_3'), methods=['POST'])
def webhook_bot3():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot3.process_new_updates([update])
        return 'ok', 200

    return 'Method Not Allowed', 405

def set_webhook():
    bot.remove_webhook()
    bot.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN)
    logging.info(f"🌍 تم تعيين الويب هوك على: {WEBHOOK_URL}/{BOT_TOKEN}")


    bot2.remove_webhook()
    bot2.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN_2)
    logging.info(f"🌍 تم تعيين الويب هوك على: {WEBHOOK_URL}/{BOT_TOKEN_2}")
    
    bot3.remove_webhook()
    bot3.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN_3)
    logging.info(f"🌍 تم تعيين الويب هوك على: {WEBHOOK_URL}/{BOT_TOKEN_3}")



import schedule
import time
import threading
from datetime import datetime

# استدعاء الدوال مباشرة إذا كانت في نفس الملف
# إذا كانت في ملف ثاني استبدل بـ from stats import ...

def run_reports():
    """تشغيل تقارير اليوم"""
    send_daily_report()
    send_top_users_report(5)  # نرسل أفضل 5

# جدولة التقارير لتعمل يومياً الساعة 23:59
schedule.every().day.at("23:59").do(run_reports)

def run_scheduler():
    """التشغيل المستمر للجدولة"""
    while True:
        schedule.run_pending()
        time.sleep(1)

# تشغيل الـ scheduler في thread منفصل
def start_scheduler():
    t = threading.Thread(target=run_scheduler, daemon=True)
    t.start()
    print(f"📅 Scheduler started at {datetime.now()}")


start_scheduler()

if __name__ == "__main__":
    init_all_dbs()
    insert_sample_quiz_if_not_exists()
    set_webhook()
    port = int(os.environ.get('PORT', 10000))  # Render يستخدم 10000
    app.run(host='0.0.0.0', port=port)

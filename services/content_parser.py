# content_parser.py
def extract_text(file_path):
    # مؤقت: هنا ضع كود استخراج النصوص من PDF/DOCX/TXT
    return "محتوى الملف كنص"
    
    
    
    
# text_extraction.py

import os
from pptx import Presentation
import docx
import fitz                     # PyMuPDF


# قائمة الملفات لكل مستخدم (تغييرها إلى dict)
user_files = {}


MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 MB

def is_file_size_allowed(bot, file_id):
    file_info = bot.get_file(file_id)
    return file_info.file_size <= MAX_FILE_SIZE



# دوال مساعدة
def is_text_empty(text):
    return not text.strip()

def save_file(uid, file_name, file_data):
    os.makedirs("downloads", exist_ok=True)
    path = os.path.join("downloads", file_name)
    with open(path, "wb") as f:
        f.write(file_data)
    user_files[uid] = path
    return path

# الدالة الرئيسية
def extract_text_from_file(uid, bot, msg, path, chat_id=None, message_id=None):
    """
    دالة استخراج النص من أي ملف.
    - uid: معرف المستخدم
    - bot: كائن البوت
    - msg: رسالة التليغرام
    - path: مسار الملف
    - chat_id, message_id: لأغراض التعديل على الرسائل أثناء الـ OCR
    """
    ext = path.rsplit(".", 1)[-1].lower()

    # تعيين دوال الاستخراج لكل نوع ملف
    extractor_map = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "txt": extract_text_from_txt,
        "pptx": extract_text_from_pptx
    }

    ocr_map = {
        "pdf": extract_text_from_pdf_with_ocr,
        "docx": extract_text_from_docx_with_ocr,
        "pptx": extract_text_from_pptx_with_ocr
    }

    if ext in ("jpg", "png"):
        if not can_generate(uid):
            return bot.send_message(uid, "⚠️ هذه الميزة متاحة فقط للمشتركين.")
        bot.edit_message_text("⏳ جاري تحليل الصورة...", chat_id=chat_id, message_id=message_id)
        content, ocr_debug = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng")
        return content

    # الحصول على دالة الاستخراج
    extract_func = extractor_map.get(ext)
    if not extract_func:
        return bot.send_message(uid, f"❌ نوع الملف '{ext}' غير مدعوم.")

    content_full = extract_func(path)
    full_length = len(content_full)

    # اقتطاع النص للمستخدمين غير المشتركين
    if not can_generate(uid):
        content = content_full[:3000]
        coverage_ratio = (len(content) / full_length) * 100 if full_length else 0
        coverage = f"{coverage_ratio:.1f}% من الملف"
    else:
        content = content_full
        coverage = "كاملة ✅"

    # إذا النص فارغ، استخدم OCR
    if is_text_empty(content):
        if not can_generate(uid):
            return bot.send_message(uid, "⚠️ لا يمكن قراءة هذا الملف تلقائيًا. تحتاج الاشتراك.")
        bot.edit_message_text("⏳ يتم تجهيز الملف... الرجاء الانتظار لحظات.", chat_id=chat_id, message_id=message_id)
        language = detect_language_from_filename(msg.document.file_name)
        ocr_func = ocr_map.get(ext)
        if ocr_func:
            content = ocr_func(path, api_key=OCR_API_KEY, language=language)

    return content, coverage

# -------------------------------------------------------------------
#                     Text Extraction & OCR
# -------------------------------------------------------------------
def extract_text_from_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        text = "\n".join([page.get_text() for page in doc])
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting PDF text: {e}")
        return ""
    # fallback to PyMuPDF text extraction
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])
    
    
    
# أضف هذه الدالة في قسم Text Extraction & OCR
def extract_text_from_docx(path: str) -> str:
    try:
        doc = docx.Document(path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting DOCX text: {e}")
        return ""




# ويجب أيضاً تعريف دالة لملفات txt
def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logging.error(f"Error extracting TXT text: {e}")
        return ""
        
def is_text_empty(text: str) -> bool:
    return not text or len(text.strip()) < 30  # يمكن تعديل الحد حسب تجربتك



def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text).strip()
    except Exception as e:
        logging.error(f"Error extracting PPTX text: {e}")
        return ""





# ---- OCR Space Integration ----
def extract_text_with_ocr_space(file_path: str, api_key="helloworld", language="eng") -> tuple:
    """
    Uses OCR.Space API to extract text from an image or scanned PDF.
    Returns: (text, debug_info)
    """
    url = 'https://api.ocr.space/parse/image'
    with open(file_path, 'rb') as f:
        response = requests.post(
            url,
            files={"file": f},
            data={
                "apikey": api_key,
                "language": language,
                "isOverlayRequired": False,
                "OCREngine": 2
            },
        )

    try:
        result = response.json()
        if result.get("IsErroredOnProcessing"):
            error_msg = result.get("ErrorMessage", "Unknown OCR error")
            return "", f"[OCR ERROR] {error_msg}"
        
        parsed = result.get("ParsedResults")
        if not parsed:
            return "", "[OCR ERROR] No ParsedResults returned."

        text = parsed[0].get("ParsedText", "").strip()
        return text, f"[OCR DEBUG] Length: {len(text)} | Excerpt: {text[:100]}"
    
    except Exception as e:
        return "", f"[OCR EXCEPTION] {e}"


# ---- PDF Split + OCR ----
def extract_text_from_pdf_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Splits a PDF into chunks of 3 pages (OCR.Space free limit),
    sends each chunk separately, and concatenates the extracted text.
    """
    try:
        doc = fitz.open(path)
        all_text = []
        # تقسيم كل 3 صفحات في ملف مؤقت
        for i in range(0, len(doc), 3):
            subdoc = fitz.open()  # ملف جديد مؤقت
            for j in range(i, min(i+3, len(doc))):
                subdoc.insert_pdf(doc, from_page=j, to_page=j)
            
            with NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                temp_path = tmp.name
                subdoc.save(temp_path)
                subdoc.close()
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)
            
            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PDF with OCR: {e}")
        return ""


# ---- PPTX Split + OCR ----
def extract_text_from_pptx_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Converts PPTX slides into smaller chunks (3 slides per file),
    sends each chunk separately to OCR.Space, and concatenates the text.
    """
    try:
        prs = Presentation(path)
        all_text = []

        # تقسيم العرض كل 3 شرائح
        for i in range(0, len(prs.slides), 3):
            new_ppt = Presentation()
            # إضافة تخطيط فارغ (مطلوب لعمل نسخ الشرائح)
            blank_layout = new_ppt.slide_layouts[6]

            for j in range(i, min(i+3, len(prs.slides))):
                slide = prs.slides[j]
                new_slide = new_ppt.slides.add_slide(blank_layout)
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        textbox = new_slide.shapes.add_textbox(left=0, top=0, width=new_ppt.slide_width, height=100)
                        textbox.text = shape.text
            
            with NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                temp_path = tmp.name
                new_ppt.save(temp_path)
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR PPTX chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)

            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PPTX with OCR: {e}")
        return ""